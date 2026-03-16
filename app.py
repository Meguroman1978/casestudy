import os
import traceback
import logging
import requests
import re
import json
import io
from urllib.parse import urlparse
from flask import Flask, render_template, request, jsonify, send_file, send_from_directory
import pandas as pd
import gspread
from google.oauth2.service_account import Credentials
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill
from concurrent.futures import ThreadPoolExecutor, as_completed
from dotenv import load_dotenv

# 環境変数をロード
load_dotenv()

# ロギング設定
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

# 国名→地域マッピング辞書
# データセット内では地域名（Americas, Europe, Japan, China/ANZ, SEA/SA/MEA, System）が使用されている
COUNTRY_TO_REGION_MAPPING = {
    # 地域名→地域名（Google Sheetの Account: Owner Territory に地域名が直接入っている場合）
    'Americas': ['Americas'],
    'Europe': ['Europe'],
    'Japan': ['Japan'],
    'China/ANZ': ['China/ANZ'],
    'SEA/SA/MEA': ['SEA/SA/MEA'],
    
    # 以下は後方互換のため保持（個別の国名が使われる場合）
    # Americas（南北アメリカ大陸）
    'United States': ['Americas'],
    'Brazil': ['Americas'],
    'Mexico': ['Americas'],
    'Canada': ['Americas'],
    'Colombia': ['Americas'],
    'Chile': ['Americas'],
    
    # Europe（ヨーロッパ）
    'Germany': ['Europe'],
    'France': ['Europe'],
    'United Kingdom': ['Europe'],
    'Italy': ['Europe'],
    'Spain': ['Europe'],
    'Poland': ['Europe'],
    'Ukraine': ['Europe'],
    'Netherlands': ['Europe'],
    'Belgium': ['Europe'],
    'Sweden': ['Europe'],
    'Austria': ['Europe'],
    'Switzerland': ['Europe'],
    'Denmark': ['Europe'],
    'Norway': ['Europe'],
    'Ireland': ['Europe'],
    'Lithuania': ['Europe'],
    
    # China/ANZ（中国・オーストラリア・ニュージーランド）
    'China': ['China/ANZ'],
    'Australia': ['China/ANZ'],
    'New Zealand': ['China/ANZ'],
    'Hong Kong': ['China/ANZ'],
    'Taiwan': ['China/ANZ'],
    
    # SEA/SA/MEA（東南アジア・南アジア・中東・アフリカ）
    'India': ['SEA/SA/MEA'],
    'Pakistan': ['SEA/SA/MEA'],
    'Thailand': ['SEA/SA/MEA'],
    'Malaysia': ['SEA/SA/MEA'],
    'Singapore': ['SEA/SA/MEA'],
    'South Korea': ['SEA/SA/MEA'],
    'Egypt': ['SEA/SA/MEA'],
    'Turkey': ['SEA/SA/MEA'],
    'South Africa': ['SEA/SA/MEA'],
    'Jordan': ['SEA/SA/MEA'],
    'United Arab Emirates': ['SEA/SA/MEA'],
    'Israel': ['SEA/SA/MEA'],
    'Qatar': ['SEA/SA/MEA']
}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.config['ALLOWED_EXTENSIONS'] = {'xlsx', 'xls'}

# uploadsディレクトリを確実に作成（Gunicorn起動時にも実行される）
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Google Sheets設定
GOOGLE_SHEET_ID = os.environ.get('GOOGLE_SHEET_ID', '')
if not GOOGLE_SHEET_ID:
    logger.warning("⚠️ GOOGLE_SHEET_ID not set in environment variables")
SCOPES = ['https://www.googleapis.com/auth/spreadsheets.readonly']

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def get_google_sheet_data():
    """Google Sheetからデータを取得する"""
    try:
        logger.info("[STEP 1] Google Sheetからデータを取得中...")
        # 公開されているシートの場合、認証なしで読み取り可能
        # CSVエクスポートURLを使用
        sheet_id = GOOGLE_SHEET_ID
        if not sheet_id:
            raise ValueError("Google Sheet ID is not configured. Please set GOOGLE_SHEET_ID environment variable.")
        gid = '0'
        csv_url = f'https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}'
        
        df = pd.read_csv(csv_url)
        logger.info(f"[STEP 1 完了] Google Sheet取得成功: {len(df)}行")
        logger.debug(f"Google Sheet カラム: {df.columns.tolist()}")
        
        # カラム名の正規化（'Business ID' -> 'Business Id'）
        if 'Business ID' in df.columns:
            df = df.rename(columns={'Business ID': 'Business Id'})
            logger.info("カラム名を正規化: 'Business ID' -> 'Business Id'")
        
        logger.debug(f"Business Id データ型: {df['Business Id'].dtype}")
        logger.debug(f"Business Id サンプル: {df['Business Id'].head(3).tolist()}")
        
        # Business Idを数値型に変換
        df['Business Id'] = pd.to_numeric(df['Business Id'], errors='coerce')
        logger.info(f"Business Idを数値型に変換: {df['Business Id'].dtype}")
        
        # Account: Industryが空欄（NaN）の場合は「不明 / Unknown」として扱う
        if 'Account: Industry' in df.columns:
            df['Account: Industry'] = df['Account: Industry'].fillna('不明 / Unknown')
            logger.info(f"空欄のAccount: Industryを「不明 / Unknown」に変換")
        
        return df
    except Exception as e:
        logger.error(f"[エラー] Google Sheet取得失敗: {e}")
        logger.error(traceback.format_exc())
        return None

def get_country_regions(country_name):
    """
    国名から対応する地域名のリストを取得
    
    Args:
        country_name: フルネームの国名（例: 'Japan', 'United States'）
    
    Returns:
        対応する地域名のリスト（例: ['Americas'], ['Japan'], ['Europe']）
        マッピングにない場合は元の国名を返す
    """
    if country_name in COUNTRY_TO_REGION_MAPPING:
        return COUNTRY_TO_REGION_MAPPING[country_name]
    else:
        # マッピングにない場合は元の国名をそのまま返す
        # （System などの特殊な値に対応）
        return [country_name]

def merge_data(video_df, live_df, sheet_df, case_type, industry_filter, country, format_filter='none', business_id_filter=''):
    """データをマージしてフィルタリングする"""
    try:
        logger.info("[STEP 2] データマージ処理開始")
        logger.info(f"選択された事例タイプ: {case_type}, 業界フィルター: {industry_filter}, 国: {country}")
        
        # 事例タイプに応じて使用するデータフレームを選択
        if case_type == 'short_video':
            if video_df is None:
                logger.error("ショート動画データが選択されていますが、ファイルがアップロードされていません")
                return None
            main_df = video_df.copy()
            logger.info("ショート動画データを使用")
        else:  # live_stream
            if live_df is None:
                logger.error("ライブ配信データが選択されていますが、ファイルがアップロードされていません")
                return None
            main_df = live_df.copy()
            logger.info("ライブ配信データを使用")
        
        logger.debug(f"選択データ: {len(main_df)}行")
        logger.debug(f"選択データのカラム: {main_df.columns.tolist()}")
        
        # 必要なカラムが存在するか確認
        if 'Business Id' not in main_df.columns:
            logger.error(f"エラー: 'Business Id' カラムが見つかりません。利用可能なカラム: {main_df.columns.tolist()}")
            return None
        
        logger.debug(f"選択データ Business Id データ型: {main_df['Business Id'].dtype}")
        logger.debug(f"選択データ Business Id サンプル: {main_df['Business Id'].head(3).tolist()}")
        logger.debug(f"Google Sheet Business Id データ型: {sheet_df['Business Id'].dtype}")
        logger.debug(f"Google Sheet Business Id サンプル: {sheet_df['Business Id'].head(3).tolist()}")
        
        # Business Idのデータ型を統一（両方を数値型に）
        main_df['Business Id'] = pd.to_numeric(main_df['Business Id'], errors='coerce')
        sheet_df['Business Id'] = pd.to_numeric(sheet_df['Business Id'], errors='coerce')
        
        logger.info("Business Idのデータ型を統一完了")
        
        # Business Idをキーとしてマージ（Channel Name, Business Nameも含める）
        logger.info("[STEP 3] データマージ実行中...")
        
        # Google Sheetに必要な列があるか確認
        available_cols = ['Business Id', 'Account: Account Name', 'Account: Industry', 'Account: Owner Territory']
        if 'Channel Name' in sheet_df.columns:
            available_cols.append('Channel Name')
        if 'Business Name' in sheet_df.columns:
            available_cols.append('Business Name')
        
        merged_df = main_df.merge(
            sheet_df[available_cols],
            on='Business Id',
            how='left'
        )
        logger.info(f"[STEP 3 完了] マージ完了: {len(merged_df)}行")
        
        # マージ結果の確認
        matched_count = merged_df['Account: Account Name'].notna().sum()
        logger.info(f"マッチング成功: {matched_count}/{len(merged_df)}行")
        
        # フィルタリング
        logger.info("[STEP 4] フィルタリング実行中...")
        before_filter = len(merged_df)
        
        # 業界フィルター（複数の業界名を受け取る場合に対応）
        if industry_filter and industry_filter != 'none':
            # カンマ区切りで複数の業界名が来る場合に対応
            if isinstance(industry_filter, str):
                industries = [i.strip() for i in industry_filter.split(',') if i.strip()]
            elif isinstance(industry_filter, list):
                industries = industry_filter
            else:
                industries = [industry_filter]
            
            if industries:
                merged_df = merged_df[merged_df['Account: Industry'].isin(industries)]
                logger.info(f"業界フィルター適用 ({', '.join(industries)}): {before_filter}行 → {len(merged_df)}行")
                before_filter = len(merged_df)
        
        # Business IDフィルター
        if business_id_filter:
            try:
                business_id_numeric = int(business_id_filter)
                logger.info(f"Business IDフィルター適用前: {before_filter}行")
                logger.info(f"検索するBusiness ID: {business_id_numeric} (型: {type(business_id_numeric)})")
                
                # フィルター適用前にBusiness Idの値をログ出力
                logger.debug(f"merged_df['Business Id']のユニーク値の数: {merged_df['Business Id'].nunique()}")
                logger.debug(f"merged_df['Business Id']のサンプル: {merged_df['Business Id'].head(10).tolist()}")
                
                # Business IDでフィルタリング
                merged_df = merged_df[merged_df['Business Id'] == business_id_numeric]
                logger.info(f"Business IDフィルター適用後: {len(merged_df)}行")
                
                if len(merged_df) == 0:
                    logger.warning(f"⚠️ Business ID {business_id_numeric} に一致する行が見つかりませんでした")
                    logger.warning(f"アップロードファイルにこのBusiness IDが含まれているか確認してください")
                else:
                    logger.info(f"✅ Business ID {business_id_numeric} で {len(merged_df)} 行が見つかりました")
                
                before_filter = len(merged_df)
            except ValueError:
                logger.warning(f"無効なBusiness ID形式: {business_id_filter}")
        
        if country != 'none':
            # 国名から対応する地域名のリストを取得
            regions = get_country_regions(country)
            logger.info(f"国名 '{country}' に対応する地域: {regions}")
            
            # データセット内の地域が、regionsのいずれかと一致する行を抽出
            merged_df = merged_df[merged_df['Account: Owner Territory'].isin(regions)]
            logger.info(f"国フィルター適用 ({country} → {regions}): {before_filter}行 → {len(merged_df)}行")
        
        logger.info(f"[STEP 4 完了] フィルタリング完了: {len(merged_df)}行")
        
        # デバッグ: フィルタリング後のデータを確認
        if len(merged_df) > 0:
            logger.debug(f"フィルタ後のBusiness Idサンプル: {merged_df['Business Id'].head(10).tolist()}")
            logger.debug(f"フィルタ後のAccount: Industryのnull数: {merged_df['Account: Industry'].isnull().sum()}/{len(merged_df)}")
        else:
            logger.warning("⚠️ フィルタリング後のデータが0行です！")
        
        # 必要な列だけを抽出
        logger.info("[STEP 5] 結果データ整形中...")
        
        # 必要な列を構築（会社名とビジネス名は削除）
        # Video Viewsカラムの名前を特定（新旧両対応）
        video_views_col = None
        if 'Video Views Uu 3 S' in merged_df.columns:
            video_views_col = 'Video Views Uu 3 S'
        elif 'VIDEO VIEWS 3S UU' in merged_df.columns:
            video_views_col = 'VIDEO VIEWS 3S UU'
        elif 'Video Views' in merged_df.columns:
            video_views_col = 'Video Views'
        
        columns_to_extract = [
            'Account: Industry',
            'Account: Owner Territory',
            'Page Url'
        ]
        
        if video_views_col:
            columns_to_extract.append(video_views_col)
        
        # 新しい指標カラムを追加（存在する場合のみ）
        # 実際のSigmaカラム名（スペース区切り）と旧カラム名の両方をサポート
        optional_metrics = [
            'View Uu Rate',  # Sigmaの実際のカラム名
            'View-Through Rate',
            'VIEW UU RATE',
            'Click Uu Rate',  # Sigmaの実際のカラム名
            'Click-Through Rate',
            'CLICK UU RATE',
            'A 2 C Uu Rate',  # Sigmaの実際のカラム名
            'Add-to-Cart Rate',
            'A2C UU RATE',
            '25 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (25%)',
            '25P VIEW COMPLETION UU RATE',
            '50 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (50%)',
            '50P VIEW COMPLETION UU RATE',
            '75 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (75%)',
            '75P VIEW COMPLETION UU RATE',
            '100 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (100%)',
            '100P VIEW_COMPLETION UU RATE'
        ]
        
        for metric in optional_metrics:
            if metric in merged_df.columns:
                columns_to_extract.append(metric)
                logger.info(f"追加指標を検出: {metric}")
        
        # Channel Nameがある場合は追加（最初に配置）
        if 'Channel Name' in merged_df.columns:
            columns_to_extract.insert(0, 'Channel Name')
        
        result_df = merged_df[columns_to_extract].copy()
        
        # 列名を日本語に変更（会社名とビジネス名は含めない）
        new_column_names = []
        if 'Channel Name' in merged_df.columns:
            new_column_names.append('チャンネル名')
        new_column_names.extend(['業種', '国', 'URL'])
        
        # Video Viewsカラムが存在する場合のみ追加
        if video_views_col:
            new_column_names.append('VIDEO_VIEWS')
        
        # 新しい指標のカラム名を追加（実際のSigmaカラム名に対応）
        if 'View Uu Rate' in merged_df.columns or 'View-Through Rate' in merged_df.columns or 'VIEW UU RATE' in merged_df.columns:
            new_column_names.append('VIEWTHROUGH_RATE')
        if 'Click Uu Rate' in merged_df.columns or 'Click-Through Rate' in merged_df.columns or 'CLICK UU RATE' in merged_df.columns:
            new_column_names.append('CLICKTHROUGH_RATE')
        if 'A 2 C Uu Rate' in merged_df.columns or 'Add-to-Cart Rate' in merged_df.columns or 'A2C UU RATE' in merged_df.columns:
            new_column_names.append('A2C_RATE')
        if '25 View Completion Uu Rate' in merged_df.columns or 'Video Completion Rate (25%)' in merged_df.columns or '25P VIEW COMPLETION UU RATE' in merged_df.columns:
            new_column_names.append('COMPLETION_RATE_25P')
        if '50 View Completion Uu Rate' in merged_df.columns or 'Video Completion Rate (50%)' in merged_df.columns or '50P VIEW COMPLETION UU RATE' in merged_df.columns:
            new_column_names.append('COMPLETION_RATE_50P')
        if '75 View Completion Uu Rate' in merged_df.columns or 'Video Completion Rate (75%)' in merged_df.columns or '75P VIEW COMPLETION UU RATE' in merged_df.columns:
            new_column_names.append('COMPLETION_RATE_75P')
        if '100 View Completion Uu Rate' in merged_df.columns or 'Video Completion Rate (100%)' in merged_df.columns or '100P VIEW_COMPLETION UU RATE' in merged_df.columns:
            new_column_names.append('COMPLETION_RATE_100P')
        
        result_df.columns = new_column_names
        
        # フォーマット列を追加（初期値は空文字列）
        result_df['フォーマット'] = ''
        
        # NaNを空文字列に変換
        result_df = result_df.fillna('')
        
        # URLからドメインを抽出
        result_df['ドメイン'] = result_df['URL'].apply(lambda x: urlparse(x).hostname if x else '')
        
        logger.info(f"[STEP 5 完了] 最終結果: {len(result_df)}行")
        logger.debug(f"結果のサンプル:\n{result_df.head(3)}")
        
        return result_df
    except Exception as e:
        logger.error(f"[エラー] データマージ失敗: {e}")
        logger.error(traceback.format_exc())
        return None

def group_by_domain_and_paginate(result_df, page=1, page_size=5):
    """Channel Name単位でグループ化してページング (デフォルト5件/ページ、各チャンネル最大3 URL)"""
    try:
        logger.info(f"[STEP 6] Channel Name単位でグループ化中... (ページ: {page}, サイズ: {page_size})")
        
        # Channel Name列が存在するか確認
        if 'チャンネル名' not in result_df.columns:
            logger.warning("チャンネル名列が存在しません。ドメインでグループ化します。")
            group_column = 'ドメイン'
        else:
            group_column = 'チャンネル名'
        
        # グループ化の列を決定（会社名とビジネス名は削除）
        agg_dict = {
            '業種': 'first',
            '国': 'first',
            'URL': 'count'
        }
        
        # VIDEO_VIEWSカラムがある場合のみ追加
        if 'VIDEO_VIEWS' in result_df.columns:
            agg_dict['VIDEO_VIEWS'] = 'sum'
        
        # オプションの指標カラムがある場合は集計に含める
        # 注: VIEWTHROUGH_RATEは削除（視聴完了率100%UU率と重複するため）
        optional_metric_cols = ['CLICKTHROUGH_RATE', 'A2C_RATE', 
                               'COMPLETION_RATE_25P', 'COMPLETION_RATE_50P', 
                               'COMPLETION_RATE_75P', 'COMPLETION_RATE_100P']
        for col in optional_metric_cols:
            if col in result_df.columns:
                # 率系は中央値、回数系は合計
                if 'RATE' in col or 'COMPLETION' in col:
                    agg_dict[col] = 'median'
                else:
                    agg_dict[col] = 'sum'
        
        # チャンネル名がある場合（グループ化対象でない場合のみ）
        if 'チャンネル名' in result_df.columns and group_column != 'チャンネル名':
            agg_dict['チャンネル名'] = 'first'
        
        # グループ化して集計
        channel_summary = result_df.groupby(group_column).agg(agg_dict).reset_index()
        
        # 合計視聴回数で降順ソート（VIDEO_VIEWSがある場合のみ）
        if 'VIDEO_VIEWS' in channel_summary.columns:
            channel_summary = channel_summary.sort_values('VIDEO_VIEWS', ascending=False)
        else:
            # VIDEO_VIEWSがない場合はURL数でソート
            channel_summary = channel_summary.sort_values('URL', ascending=False)
        
        logger.info(f"グループ化完了: Top {len(channel_summary)}件")
        
        # ページネーション適用
        start_idx = (page - 1) * page_size
        end_idx = start_idx + page_size
        paginated_channels = channel_summary.iloc[start_idx:end_idx]
        
        logger.info(f"ページ {page}: {len(paginated_channels)}件を返却")
        
        # ページネーション対象のチャンネルの詳細データを取得
        channel_list = paginated_channels[group_column].tolist()
        detailed_data = result_df[result_df[group_column].isin(channel_list)].copy()
        
        # 視聴回数で降順ソート（チャンネル内）- VIDEO_VIEWSがある場合のみ
        if 'VIDEO_VIEWS' in detailed_data.columns:
            detailed_data = detailed_data.sort_values([group_column, 'VIDEO_VIEWS'], ascending=[True, False])
        else:
            detailed_data = detailed_data.sort_values(group_column, ascending=True)
        
        # 各チャンネルのURL数を最大3に制限
        detailed_data = detailed_data.groupby(group_column).head(3).reset_index(drop=True)
        
        return {
            'channel_summary': channel_summary,
            'paginated_channels': paginated_channels,
            'detailed_data': detailed_data,
            'total_domains': len(channel_summary),  # Changed from total_channels to total_domains
            'current_page': page,
            'page_size': page_size,
            'has_next': end_idx < len(channel_summary)
        }
    except Exception as e:
        logger.error(f"[エラー] ドメイングループ化失敗: {e}")
        logger.error(traceback.format_exc())
        return None

@app.route('/')
def index():
    """トップページ"""
    return render_template('index.html')

@app.route('/favicon.ico')
def favicon():
    """Favicon"""
    return send_from_directory(os.path.join(app.root_path, 'static'),
                               'favicon.ico', mimetype='image/vnd.microsoft.icon')

@app.route('/api/get-options', methods=['GET'])
def get_options():
    """業界名と国のオプションを取得"""
    try:
        sheet_df = get_google_sheet_data()
        
        # 地域リスト（データセットの実際の地域分類に基づく）
        fixed_regions = [
            'Americas',      # 南北アメリカ大陸（US, Brazil, Mexico, Canada, Colombia, Chile）
            'Europe',        # ヨーロッパ（Germany, France, UK, Italy, Spain, etc.）
            'Japan',         # 日本
            'China/ANZ',     # 中国・オーストラリア・ニュージーランド
            'SEA/SA/MEA'     # 東南アジア・南アジア・中東・アフリカ
        ]
        
        # Google Sheetからデータを取得できない場合
        if sheet_df is None:
            return jsonify({
                'industries': [],
                'countries': fixed_regions
            })
        
        # ユニークな業界名を取得（空でないもの）
        industries = sorted(sheet_df['Account: Industry'].dropna().unique().tolist())
        
        # 地域リストを使用（データセットの実際の地域分類に基づく）
        return jsonify({
            'industries': industries,
            'countries': fixed_regions
        })
    except Exception as e:
        # エラーが発生した場合も地域リストを返す
        fixed_regions = [
            'Americas',      # 南北アメリカ大陸
            'Europe',        # ヨーロッパ
            'Japan',         # 日本
            'China/ANZ',     # 中国・オーストラリア・ニュージーランド
            'SEA/SA/MEA'     # 東南アジア・南アジア・中東・アフリカ
        ]
        return jsonify({
            'industries': [],
            'countries': fixed_regions
        })

@app.route('/api/get-category-hierarchy', methods=['GET'])
def get_category_hierarchy():
    """カテゴリー階層を取得"""
    try:
        hierarchy_path = os.path.join(os.path.dirname(__file__), 'category_hierarchy.json')
        with open(hierarchy_path, 'r', encoding='utf-8') as f:
            hierarchy = json.load(f)
        return jsonify(hierarchy)
    except Exception as e:
        logger.error(f"Category hierarchy load error: {e}")
        return jsonify({'error': str(e)}), 500

def detect_firework_format(html_content):
    """
    HTMLコンテンツからFireworkスクリプトのフォーマットを検出
    
    Returns:
        str: 検出されたフォーマット名（複数ある場合は最初のもの）、なければ'Unknown'
    """
    if not html_content:
        return 'Unknown'
    
    # 各フォーマットの検出パターン（優先順位順）
    format_patterns = [
        # Horizontal Carousel: style属性にthumbnailが含まれる
        (r'<fw-embed-feed[^>]*style=["\'][^"\']*thumbnail[^"\']*["\'][^>]*>', 'Horizontal Carousel'),
        
        # Dynamic Carousel: thumbnail_style="dynamic"
        (r'<fw-embed-feed[^>]*thumbnail_style=["\']dynamic["\'][^>]*>', 'Dynamic Carousel'),
        
        # Grid: mode="grid"
        (r'<fw-embed-feed[^>]*mode=["\']grid["\'][^>]*>', 'Grid'),
        
        # Carousel: mode="row"
        (r'<fw-embed-feed[^>]*mode=["\']row["\'][^>]*>', 'Carousel'),
        
        # Story Block: <fw-embed-feed> without mode/style/thumbnail_style
        (r'<fw-embed-feed(?![^>]*(?:mode=|style=|thumbnail_style=))[^>]*>', 'Story Block'),
        
        # Circle Stories: thumbnail_shape="circle"
        (r'<fw-stories[^>]*thumbnail_shape=["\']circle["\'][^>]*>', 'Circle Stories'),
        
        # Vertical Stories: thumbnail_shape="rectangle"
        (r'<fw-stories[^>]*thumbnail_shape=["\']rectangle["\'][^>]*>', 'Vertical Stories'),
        
        # Floating Player: mode="pinned"
        (r'<fw-storyblock[^>]*mode=["\']pinned["\'][^>]*>', 'Floating Player'),
        
        # Horizontal Player
        (r'<fw-player[^>]*>', 'Horizontal Player'),
        
        # Hero Unit
        (r'<fw-herounit[^>]*>', 'Hero Unit'),
        
        # Player Deck
        (r'<fw-player-deck[^>]*>', 'Player Deck'),
    ]
    
    # 各パターンをチェック
    for pattern, format_name in format_patterns:
        if re.search(pattern, html_content, re.IGNORECASE | re.DOTALL):
            logger.debug(f"Detected format: {format_name}")
            return format_name
    
    # どのパターンにもマッチしない場合
    if re.search(r'<fw-[\w-]+', html_content, re.IGNORECASE):
        logger.debug("Firework tag found but format unknown")
        return 'Unknown'
    
    return 'Unknown'

def check_fw_tag_in_url(url):
    """指定されたURLのソースコードに<fw-タグが含まれているかチェックし、フォーマットも検出"""
    try:
        logger.info(f"Checking <fw- tag for URL: {url}")
        # タイムアウトを設定してページを取得
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        html_content = response.text
        
        # <fw- で始まるタグを検索
        has_fw_tag = bool(re.search(r'<fw-[\w-]+', html_content, re.IGNORECASE))
        logger.info(f"<fw- tag found: {has_fw_tag}")
        
        # フォーマットを検出
        format_name = detect_firework_format(html_content) if has_fw_tag else 'Unknown'
        logger.info(f"Detected format: {format_name}")
        
        return has_fw_tag, html_content, format_name
    except Exception as e:
        logger.error(f"Error checking <fw- tag: {e}")
        return False, None, 'Unknown'

@app.route('/api/check-fw-tag', methods=['GET'])
def api_check_fw_tag():
    """URLの<fw-タグチェックとスクリーンショット情報を返すAPI"""
    try:
        url = request.args.get('url')
        if not url:
            return jsonify({'error': 'URL parameter is required'}), 400
        
        has_fw_tag, html_content = check_fw_tag_in_url(url)
        
        # スクリーンショットURL（要件5用）
        # 複数のスクリーンショットサービスを試行
        screenshot_url = None
        if has_fw_tag:
            # URLエンコード
            from urllib.parse import quote
            encoded_url = quote(url, safe='')
            
            # Option 1: screenshotapi.net (無料、登録不要)
            screenshot_url = f"https://shot.screenshotapi.net/screenshot?url={encoded_url}&width=400&height=300&output=image&file_type=png&wait_for_event=load"
            
            # Option 2 (fallback): Google's Page Speed Insights screenshot
            # screenshot_url = f"https://www.googleapis.com/pagespeedonline/v5/runPagespeed?url={encoded_url}&screenshot=true"
            
            # Option 3 (fallback): screenshot.guru
            # screenshot_url = f"https://api.screenshot.guru/screenshot?url={encoded_url}&width=400&height=300"
        
        return jsonify({
            'has_fw_tag': has_fw_tag,
            'screenshot_url': screenshot_url
        })
    except Exception as e:
        logger.error(f"Error in check_fw_tag API: {e}")
        logger.error(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

@app.route('/api/process', methods=['POST'])
def process_data():
    """アップロードされたファイルを処理"""
    try:
        logger.info("="*60)
        logger.info("新しい検索リクエスト開始")
        logger.info("="*60)
        
        # ファイルのチェック（少なくとも1つのファイルが必要）
        video_file = request.files.get('video_file')
        live_file = request.files.get('live_file')
        
        has_video = video_file and video_file.filename != ''
        has_live = live_file and live_file.filename != ''
        
        if not has_video and not has_live:
            logger.warning("ファイルがアップロードされていません")
            return jsonify({'error': '少なくとも1つのファイルをアップロードしてください'}), 400
        
        # アップロードされたファイルの形式チェック
        if has_video and not allowed_file(video_file.filename):
            logger.warning(f"不正なファイル形式: {video_file.filename}")
            return jsonify({'error': 'Excelファイル (.xlsx, .xls) のみアップロード可能です'}), 400
        
        if has_live and not allowed_file(live_file.filename):
            logger.warning(f"不正なファイル形式: {live_file.filename}")
            return jsonify({'error': 'Excelファイル (.xlsx, .xls) のみアップロード可能です'}), 400
        
        # パラメータの取得
        case_type = request.form.get('case_type', 'short_video')
        industry_filter = request.form.get('industry_filter', 'none')
        country = request.form.get('country', 'none')
        format_filter = request.form.get('format_filter', 'none')
        business_id_filter = request.form.get('business_id_filter', '').strip()
        page = int(request.form.get('page', 1))
        page_size = int(request.form.get('page_size', 5))  # デフォルトを5に変更（最大5チャンネル表示）
        
        logger.info(f"検索条件: 事例タイプ={case_type}, 業界フィルター={industry_filter}, 国={country}, フォーマット={format_filter}, Business ID={business_id_filter}, ページ={page}")
        logger.info(f"アップロードされたファイル: video={has_video}, live={has_live}")
        
        # uploadsディレクトリの存在を確認（念のため）
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        
        # ファイルを一時保存（アップロードされたもののみ）
        video_path = None
        live_path = None
        
        if has_video:
            video_filename = secure_filename(video_file.filename)
            video_path = os.path.join(app.config['UPLOAD_FOLDER'], video_filename)
            video_file.save(video_path)
            video_size = os.path.getsize(video_path) / (1024 * 1024)  # MB
            logger.info(f"ショート動画ファイル保存完了: {video_filename} ({video_size:.2f}MB)")
            
            if video_size > 10:
                logger.warning(f"⚠️ 大きなファイルが検出されました: video={video_size:.2f}MB")
                logger.warning(f"⚠️ 処理に時間がかかる可能性があります（最大5分）")
        
        if has_live:
            live_filename = secure_filename(live_file.filename)
            live_path = os.path.join(app.config['UPLOAD_FOLDER'], live_filename)
            live_file.save(live_path)
            live_size = os.path.getsize(live_path) / (1024 * 1024)  # MB
            logger.info(f"ライブ配信ファイル保存完了: {live_filename} ({live_size:.2f}MB)")
            
            if live_size > 10:
                logger.warning(f"⚠️ 大きなファイルが検出されました: live={live_size:.2f}MB")
                logger.warning(f"⚠️ 処理に時間がかかる可能性があります（最大5分）")
        
        # データの読み込み（メモリ効率化）
        logger.info("[STEP 0] Excelファイル読み込み中...")
        
        # 必要なカラムのみ読み込んでメモリを節約（新しい指標を追加）
        required_columns = ['Page Url', 'Business Id', 'Business Name', 'Business Country', 
                          'Channel Id', 'Channel Name']
        
        # Video Viewsカラムの名前を候補に追加（新旧両対応）
        video_views_candidates = ['Video Views Uu 3 S', 'VIDEO VIEWS 3S UU', 'Video Views']
        
        # オプションの指標カラム（存在する場合のみ読み込む）- 実際のSigmaカラム名を追加
        optional_metrics = [
            'View Uu Rate',  # Sigmaの実際のカラム名
            'View-Through Rate',
            'VIEW UU RATE',
            'Click Uu Rate',  # Sigmaの実際のカラム名
            'Click-Through Rate',
            'CLICK UU RATE',
            'A 2 C Uu Rate',  # Sigmaの実際のカラム名
            'Add-to-Cart Rate',
            'A2C UU RATE',
            '25 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (25%)',
            '25P VIEW COMPLETION UU RATE',
            '50 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (50%)',
            '50P VIEW COMPLETION UU RATE',
            '75 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (75%)',
            '75P VIEW COMPLETION UU RATE',
            '100 View Completion Uu Rate',  # Sigmaの実際のカラム名
            'Video Completion Rate (100%)',
            '100P VIEW_COMPLETION UU RATE'
        ]
        
        # データの読み込み（アップロードされたファイルのみ）
        logger.info("[STEP 0] Excelファイル読み込み中...")
        
        video_df = None
        live_df = None
        
        if has_video:
            try:
                # read_excel with engine='openpyxl' and read_only=True for memory efficiency
                video_df = pd.read_excel(video_path, engine='openpyxl')
                logger.info(f"ショート動画データ: {len(video_df)}行, カラム: {video_df.columns.tolist()}")
                
                # Business Idのデータを確認
                if 'Business Id' in video_df.columns:
                    logger.info(f"Business Idのユニーク値の数: {video_df['Business Id'].nunique()}")
                    logger.debug(f"Business Idのサンプル（最初の10個）: {video_df['Business Id'].head(10).tolist()}")
                    logger.debug(f"Business Idのサンプル（最後の10個）: {video_df['Business Id'].tail(10).tolist()}")
                
                # 不要なカラムを削除してメモリを解放（オプションの指標も含める）
                video_columns_to_keep = [col for col in required_columns if col in video_df.columns]
                
                # Video Viewsカラムを検出して追加
                for vv_col in video_views_candidates:
                    if vv_col in video_df.columns:
                        video_columns_to_keep.append(vv_col)
                        logger.info(f"Video Viewsカラムを検出: {vv_col}")
                        break
                
                # オプションの指標カラムを追加
                for metric in optional_metrics:
                    if metric in video_df.columns:
                        video_columns_to_keep.append(metric)
                        logger.info(f"オプション指標を検出: {metric}")
                
                if video_columns_to_keep:
                    video_df = video_df[video_columns_to_keep]
                    logger.info(f"必要なカラムのみ保持: {video_columns_to_keep}")
            except Exception as e:
                logger.error(f"ショート動画ファイル読み込みエラー: {e}")
                raise
        
        if has_live:
            try:
                live_df = pd.read_excel(live_path, engine='openpyxl')
                logger.info(f"ライブ配信データ: {len(live_df)}行, カラム: {live_df.columns.tolist()}")
                
                # Business Idのデータを確認
                if 'Business Id' in live_df.columns:
                    logger.info(f"Business Idのユニーク値の数: {live_df['Business Id'].nunique()}")
                    logger.debug(f"Business Idのサンプル（最初の10個）: {live_df['Business Id'].head(10).tolist()}")
                    logger.debug(f"Business Idのサンプル（最後の10個）: {live_df['Business Id'].tail(10).tolist()}")
                
                # 不要なカラムを削除してメモリを解放（オプションの指標も含める）
                live_columns_to_keep = [col for col in required_columns if col in live_df.columns]
                
                # Video Viewsカラムを検出して追加
                for vv_col in video_views_candidates:
                    if vv_col in live_df.columns:
                        live_columns_to_keep.append(vv_col)
                        logger.info(f"Video Viewsカラムを検出（ライブ）: {vv_col}")
                        break
                
                # オプションの指標カラムを追加
                for metric in optional_metrics:
                    if metric in live_df.columns:
                        live_columns_to_keep.append(metric)
                        logger.info(f"オプション指標を検出（ライブ）: {metric}")
                
                if live_columns_to_keep:
                    live_df = live_df[live_columns_to_keep]
                    logger.info(f"必要なカラムのみ保持: {live_columns_to_keep}")
            except Exception as e:
                logger.error(f"ライブ配信ファイル読み込みエラー: {e}")
                raise
        
        sheet_df = get_google_sheet_data()
        
        if sheet_df is None:
            logger.error("Google Sheetデータの取得に失敗")
            return jsonify({'error': 'Google Sheetからデータを取得できませんでした'}), 500
        
        # データのマージとフィルタリング
        result_df = merge_data(video_df, live_df, sheet_df, case_type, industry_filter, country, format_filter, business_id_filter)
        
        if result_df is None:
            logger.error("データマージ処理に失敗")
            # より具体的なエラーメッセージを返す
            if case_type == 'short_video' and video_df is None:
                error_msg = 'ショート動画事例を検索するには、Top Video Views per URL ファイルをアップロードしてください'
            elif case_type == 'live_stream' and live_df is None:
                error_msg = 'ライブ配信事例を検索するには、Top Live Stream Views per URL ファイルをアップロードしてください'
            else:
                error_msg = 'データの処理中にエラーが発生しました。アップロードしたファイルに必要なカラム（Business Id, Page Url等）が含まれているか確認してください。'
            return jsonify({'error': error_msg}), 500
        
        # ドメインごとにグループ化してページネーション
        pagination_result = group_by_domain_and_paginate(result_df, page=page, page_size=page_size)
        
        if pagination_result is None:
            logger.error("ドメインのグループ化に失敗")
            return jsonify({'error': 'ドメインのグループ化中にエラーが発生しました。'}), 500
        
        # 要件4: 表示対象のURLのみ<fw-タグをチェック（パフォーマンス改善）
        logger.info("[STEP 7] 表示対象URLの<fw-タグチェック開始...")
        detailed_data = pagination_result['detailed_data']
        original_count = len(detailed_data)
        
        # 並列処理でURLチェックを高速化 (最大10スレッド並列)
        logger.info(f"並列URLチェック開始: {original_count}件のURL")
        
        def check_url_wrapper(url):
            """URLチェックのラッパー関数（フォーマット検出も含む）"""
            try:
                has_fw_tag, html_content, format_name = check_fw_tag_in_url(url)
                return (url, has_fw_tag, format_name)
            except Exception as e:
                logger.error(f"URLチェック失敗 ({url}): {e}")
                return (url, False, 'Unknown')
        
        # 並列実行
        url_to_flag = {}
        url_to_format = {}
        with ThreadPoolExecutor(max_workers=10) as executor:
            # すべてのURLに対してチェックを送信
            future_to_url = {executor.submit(check_url_wrapper, row['URL']): row['URL'] 
                           for idx, row in detailed_data.iterrows()}
            
            # 完了したものから結果を取得
            completed = 0
            for future in as_completed(future_to_url):
                url, has_fw_tag, format_name = future.result()
                url_to_flag[url] = has_fw_tag
                url_to_format[url] = format_name
                completed += 1
                if completed % 10 == 0:
                    logger.info(f"進行状況: {completed}/{original_count} URLチェック完了")
        
        # 結果をDataFrameに反映
        detailed_data['has_fw_tag'] = detailed_data['URL'].map(url_to_flag)
        detailed_data['フォーマット'] = detailed_data['URL'].map(url_to_format)
        
        # fwタグがあるものだけをフィルター
        filtered_data = detailed_data[detailed_data['has_fw_tag'] == True].copy()
        
        # フォーマットフィルター適用
        if format_filter and format_filter != 'none' and format_filter != '選択しない':
            before_format_filter = len(filtered_data)
            filtered_data = filtered_data[filtered_data['フォーマット'] == format_filter].copy()
            logger.info(f"フォーマットフィルター適用 ({format_filter}): {before_format_filter}行 → {len(filtered_data)}行")
        
        # 内部使用列を削除（フロントエンドで使用しないため）
        if 'ドメイン' in filtered_data.columns:
            filtered_data = filtered_data.drop('ドメイン', axis=1)
        if 'has_fw_tag' in filtered_data.columns:
            filtered_data = filtered_data.drop('has_fw_tag', axis=1)
        
        # 列の順序を調整：チャンネル名、業種、国、フォーマット、URL、VIDEO_VIEWS、その他の指標
        desired_order = []
        if 'チャンネル名' in filtered_data.columns:
            desired_order.append('チャンネル名')
        if '業種' in filtered_data.columns:
            desired_order.append('業種')
        if '国' in filtered_data.columns:
            desired_order.append('国')
        if 'フォーマット' in filtered_data.columns:
            desired_order.append('フォーマット')
        if 'URL' in filtered_data.columns:
            desired_order.append('URL')
        
        # 指標カラムを追加（存在する場合のみ）
        # 注: VIEWTHROUGH_RATEは削除（視聴完了率100%UU率と重複するため）
        metric_columns = ['VIDEO_VIEWS', 'CLICKTHROUGH_RATE', 'A2C_RATE',
                         'COMPLETION_RATE_25P', 'COMPLETION_RATE_50P', 
                         'COMPLETION_RATE_75P', 'COMPLETION_RATE_100P']
        for col in metric_columns:
            if col in filtered_data.columns:
                desired_order.append(col)
        
        # 順序通りに列を並べ替え
        filtered_data = filtered_data[desired_order]
        
        logger.info(f"[STEP 7 完了] <fw-タグフィルター: {original_count}行 → {len(filtered_data)}行")
        
        # 一時ファイルを削除（アップロードされたファイルのみ）
        try:
            if video_path and os.path.exists(video_path):
                os.remove(video_path)
                logger.info("video_path削除完了")
            if live_path and os.path.exists(live_path):
                os.remove(live_path)
                logger.info("live_path削除完了")
        except Exception as cleanup_error:
            logger.warning(f"一時ファイル削除エラー: {cleanup_error}")
        
        # 結果をJSON形式で返す
        result = {
            'columns': filtered_data.columns.tolist(),
            'data': filtered_data.to_dict(orient='records'),
            'total_count': len(filtered_data),
            'total_domains': pagination_result['total_domains'],
            'current_page': pagination_result['current_page'],
            'page_size': pagination_result['page_size'],
            'has_next': pagination_result['has_next']
        }
        
        logger.info("="*60)
        logger.info(f"検索成功: {len(filtered_data)}件の結果を返却 (ページ {page}/{(pagination_result['total_domains'] + page_size - 1) // page_size})")
        logger.info("="*60)
        
        return jsonify(result)
    
    except Exception as e:
        logger.error("="*60)
        logger.error(f"予期しないエラーが発生: {e}")
        logger.error(f"エラータイプ: {type(e).__name__}")
        logger.error(traceback.format_exc())
        logger.error("="*60)
        
        # エラー時も一時ファイルを削除
        try:
            if 'video_path' in locals() and os.path.exists(video_path):
                os.remove(video_path)
                logger.info("エラー時: video_path削除完了")
            if 'live_path' in locals() and os.path.exists(live_path):
                os.remove(live_path)
                logger.info("エラー時: live_path削除完了")
        except Exception as cleanup_error:
            logger.warning(f"一時ファイル削除エラー: {cleanup_error}")
        
        # より詳細なエラーメッセージ
        error_detail = str(e)
        if "No such file or directory" in error_detail:
            error_msg = 'ファイルの保存に失敗しました。サーバーの設定を確認してください。 / File save failed. Please check server configuration.'
        elif "Google Sheet" in error_detail or "gspread" in error_detail:
            error_msg = 'Google Sheetsへのアクセスに失敗しました。GOOGLE_SHEET_IDを確認してください。 / Failed to access Google Sheets. Please check GOOGLE_SHEET_ID.'
        elif "pandas" in error_detail or "read_excel" in error_detail or "openpyxl" in error_detail:
            error_msg = 'Excelファイルの読み込みに失敗しました。ファイルが大きすぎるか、形式が正しくない可能性があります。 / Failed to read Excel file. File may be too large or format is incorrect.'
        elif "timeout" in error_detail.lower() or "SIGKILL" in error_detail:
            error_msg = 'ファイルの処理がタイムアウトしました。ファイルサイズを小さくしてください（推奨: 5MB以下）。 / File processing timed out. Please reduce file size (recommended: under 5MB).'
        else:
            error_msg = f'エラーが発生しました / Error occurred: {error_detail}'
        
        return jsonify({'error': error_msg, 'detail': error_detail}), 500

def extract_website_info(url):
    """ウェブサイトからメタ情報を抽出"""
    try:
        response = requests.get(url, timeout=10, headers={'User-Agent': 'Mozilla/5.0'})
        html_content = response.text
        
        # メタタグから情報を抽出
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 会社概要（descriptionメタタグから）
        description_tag = soup.find('meta', attrs={'name': 'description'}) or soup.find('meta', attrs={'property': 'og:description'})
        description = description_tag.get('content', '') if description_tag else ''
        
        # タイトル
        title = soup.find('title').get_text() if soup.find('title') else ''
        
        # ロゴ画像URL
        logo_url = None
        logo_tag = soup.find('meta', attrs={'property': 'og:image'}) or soup.find('link', attrs={'rel': 'icon'}) or soup.find('link', attrs={'rel': 'apple-touch-icon'})
        if logo_tag:
            logo_url = logo_tag.get('content') or logo_tag.get('href')
            if logo_url and not logo_url.startswith('http'):
                from urllib.parse import urljoin
                logo_url = urljoin(url, logo_url)
        
        return {
            'title': title,
            'description': description,
            'logo_url': logo_url
        }
    except Exception as e:
        logger.error(f"ウェブサイト情報抽出エラー: {e}")
        return {'title': '', 'description': '', 'logo_url': None}

def translate_text(text, target_lang='en'):
    """テキストを翻訳（簡易版 - 実際にはGoogle Translate APIなどを使用）"""
    # ここでは簡易的に、日本語が含まれている場合のみ翻訳を試みる
    if not text or target_lang == 'ja':
        return text
    
    # 日本語が含まれているか確認
    import re
    if re.search(r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]', text):
        # 実際の実装ではGoogle Translate APIを使用
        # ここでは元のテキストを返す（翻訳APIを実装する場合は置き換え）
        return text
    
    return text

def search_logo_images(channel_name, country='', industry='', count=3):
    """Channel name + Country + Industryでロゴ画像を検索（検索精度向上）"""
    try:
        from bs4 import BeautifulSoup
        import urllib.parse
        
        # 検索クエリを構築: Channel Name + Country + Industry で検索精度向上
        search_parts = [channel_name]
        if country and country != '選択しない':
            search_parts.append(country)
        if industry and industry != '選択しない':
            search_parts.append(industry)
        search_parts.append('logo')
        
        search_query = ' '.join(search_parts)
        logger.info(f"Searching logos with query: {search_query}")
        
        # Google画像検索のURL（スクレイピング）
        encoded_query = urllib.parse.quote(search_query)
        search_url = f"https://www.google.com/search?q={encoded_query}&tbm=isch"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(search_url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # 画像URLを抽出
        logo_urls = []
        img_tags = soup.find_all('img')
        
        for img in img_tags[:count + 5]:  # 余分に取得
            img_url = img.get('src') or img.get('data-src')
            if img_url and img_url.startswith('http') and len(logo_urls) < count:
                # base64やデータURLは除外
                if not img_url.startswith('data:'):
                    logo_urls.append(img_url)
        
        logger.info(f"Found {len(logo_urls)} logo URLs")
        return logo_urls[:count]
        
    except Exception as e:
        logger.error(f"Logo search error: {e}")
        logger.error(traceback.format_exc())
        return []

def save_complete_html_page(url, output_path):
    """
    Playwrightを使用してページ全体を単一のHTMLファイルとして保存
    （CSS、画像などすべてのリソースをインライン化、リッチデザイン維持）
    
    Args:
        url: 保存対象のURL
        output_path: 保存先のファイルパス
    
    Returns:
        bool: 成功した場合True
    """
    try:
        from playwright.sync_api import sync_playwright
        
        logger.info(f"Saving complete HTML page for: {url}")
        
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            context = browser.new_context(
                viewport={'width': 1200, 'height': 800},
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            )
            page = context.new_page()
            
            # ページに移動（最も緩い条件から試す）
            try:
                logger.info("Loading page with domcontentloaded wait...")
                page.goto(url, wait_until='domcontentloaded', timeout=10000)
                logger.info("✅ Page loaded with domcontentloaded")
                # 追加の待機時間でJavaScriptが完全に実行されるのを待つ
                page.wait_for_timeout(3000)
            except Exception as e:
                logger.warning(f"domcontentloaded failed, falling back to load: {e}")
                try:
                    page.goto(url, wait_until='load', timeout=15000)
                    logger.info("✅ Page loaded with load event")
                    page.wait_for_timeout(3000)
                except Exception as e2:
                    logger.warning(f"load failed, falling back to commit: {e2}")
                    page.goto(url, wait_until='commit', timeout=10000)
                    logger.info("✅ Page loaded with commit")
                    page.wait_for_timeout(5000)  # commitの場合は長めに待機
            
            # 🎨 リッチなHTMLコンテンツを取得（より包括的なスタイル保存）
            complete_html = page.evaluate("""
                async () => {
                    // すべての画像のsrcを絶対URLに変換
                    document.querySelectorAll('img').forEach(img => {
                        if (img.src) {
                            img.setAttribute('src', img.src);
                        }
                        if (img.srcset) {
                            img.setAttribute('srcset', img.srcset);
                        }
                    });
                    
                    // すべてのリンク（CSS）のhrefを絶対URLに変換
                    document.querySelectorAll('link[href]').forEach(link => {
                        link.setAttribute('href', link.href);
                    });
                    
                    // すべてのスクリプトのsrcを絶対URLに変換
                    document.querySelectorAll('script[src]').forEach(script => {
                        script.setAttribute('src', script.src);
                    });
                    
                    // 🆕 すべての要素に計算されたスタイルをインライン化（より包括的）
                    // body配下のすべての可視要素を対象
                    const allElements = document.querySelectorAll('body *');
                    const styleProps = [
                        'background-color', 'background-image', 'background-size', 'background-position',
                        'color', 'font-family', 'font-size', 'font-weight', 'font-style',
                        'padding', 'padding-top', 'padding-right', 'padding-bottom', 'padding-left',
                        'margin', 'margin-top', 'margin-right', 'margin-bottom', 'margin-left',
                        'border', 'border-radius', 'border-width', 'border-style', 'border-color',
                        'display', 'position', 'width', 'height', 'max-width', 'max-height',
                        'flex', 'flex-direction', 'justify-content', 'align-items',
                        'text-align', 'line-height', 'letter-spacing',
                        'opacity', 'z-index', 'box-shadow', 'text-shadow',
                        'transform', 'transition'
                    ];
                    
                    let inlinedCount = 0;
                    allElements.forEach(el => {
                        try {
                            const computedStyle = window.getComputedStyle(el);
                            let inlineStyle = el.getAttribute('style') || '';
                            
                            styleProps.forEach(prop => {
                                const value = computedStyle.getPropertyValue(prop);
                                if (value && 
                                    value !== 'none' && 
                                    value !== 'normal' && 
                                    value !== 'auto' &&
                                    value !== 'rgba(0, 0, 0, 0)' &&
                                    value !== 'transparent') {
                                    inlineStyle += `${prop}:${value};`;
                                }
                            });
                            
                            if (inlineStyle) {
                                el.setAttribute('style', inlineStyle);
                                inlinedCount++;
                            }
                        } catch(e) {
                            // 個別要素のエラーは無視して続行
                        }
                    });
                    
                    console.log('Inlined styles for', inlinedCount, 'elements');
                    
                    // 🆕 <style>タグの内容も保持（既存のCSSルールを維持）
                    document.querySelectorAll('style').forEach(styleTag => {
                        styleTag.setAttribute('data-original', 'true');
                    });
                    
                    return document.documentElement.outerHTML;
                }
            """)
            
            # 🆕 外部CSSファイルをダウンロードしてインライン化
            try:
                logger.info("Downloading and inlining external CSS files...")
                
                # CSSリンクを取得
                css_links = page.evaluate("""
                    () => {
                        return Array.from(document.querySelectorAll('link[rel="stylesheet"]'))
                            .map(link => link.href)
                            .filter(href => href && href.startsWith('http'));
                    }
                """)
                
                logger.info(f"Found {len(css_links)} external CSS files")
                
                # CSSをダウンロード
                downloaded_css = []
                for css_url in css_links[:20]:  # 最大20個まで（パフォーマンス考慮）
                    try:
                        logger.info(f"Downloading CSS: {css_url}")
                        css_response = requests.get(css_url, timeout=10, headers={
                            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                        })
                        if css_response.status_code == 200:
                            downloaded_css.append(css_response.text)
                            logger.info(f"✅ Downloaded {len(css_response.text)} bytes from {css_url}")
                    except Exception as css_error:
                        logger.warning(f"Failed to download CSS {css_url}: {css_error}")
                
                # CSSを<style>タグとして追加
                if downloaded_css:
                    css_style_tag = '<style data-inlined-external="true">\n' + '\n'.join(downloaded_css) + '\n</style>'
                    # <head>タグの最後に追加
                    complete_html = complete_html.replace('</head>', f'{css_style_tag}\n</head>')
                    logger.info(f"✅ Inlined {len(downloaded_css)} external CSS files")
                
            except Exception as css_error:
                logger.warning(f"Failed to inline external CSS: {css_error}")
            
            # HTMLファイルとして保存
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(complete_html)
            
            browser.close()
            logger.info(f"✅ Rich HTML page with inlined CSS saved successfully: {output_path}")
            return True
            
    except Exception as e:
        logger.error(f"Failed to save HTML page: {e}")
        logger.error(traceback.format_exc())
        return False

def capture_screenshot_with_api(url, width=1200, height=800):
    """
    外部APIを使用してスクリーンショットを撮影（複数のAPIでフォールバック）
    
    Args:
        url: スクリーンショット対象のURL
        width: 画像幅
        height: 画像高さ
        
    Returns:
        BytesIO object containing PNG image, or None if all APIs failed
    """
    print(f"📸 capture_screenshot_with_api CALLED: url={url}")
    from urllib.parse import quote
    import requests
    
    encoded_url = quote(url, safe='')
    
    # Screenshot Machine API Key（環境変数から取得、なければ空文字列）
    screenshot_api_key = os.environ.get('SCREENSHOT_API_KEY', 'ac10e6')
    
    # 試行するAPI（優先順位順）
    apis = [
        # API 1: screenshotmachine.com（API KEY: ac10e6）
        {
            'name': 'screenshotmachine.com',
            'url': f"https://api.screenshotmachine.com/?key={screenshot_api_key}&url={encoded_url}&dimension={width}x{height}&device=desktop&format=png&cacheLimit=0&delay=1000",
            'timeout': 20
        },
        # API 2: screenshotapi.net（無料、登録不要）
        {
            'name': 'screenshotapi.net',
            'url': f"https://shot.screenshotapi.net/screenshot?url={encoded_url}&width={width}&height={height}&output=image&file_type=png&wait_for_event=load&delay=2000&full_page=false",
            'timeout': 20
        },
        # API 3: thumbnail.ws（無料、登録不要、シンプル）
        {
            'name': 'thumbnail.ws',
            'url': f"https://api.thumbnail.ws/api/{encoded_url}/viewport/{width}x{height}/fullsize",
            'timeout': 20
        },
        # API 4: Google PageSpeed Insights（無料、安定、ただしBase64デコードが必要）
        {
            'name': 'pagespeed.google',
            'url': f"https://www.googleapis.com/pagespeedonline/v5/runPagespeed?url={encoded_url}&screenshot=true",
            'timeout': 30,
            'extract_screenshot': True  # 特別処理が必要
        }
    ]
    
    # 各APIを試す
    for api in apis:
        try:
            logger.info(f"📸 Trying {api['name']}: {api['url'][:80]}...")
            print(f"📸 Trying {api['name']}...")
            
            response = requests.get(api['url'], timeout=api['timeout'], headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            })
            
            logger.info(f"📸 {api['name']} response: status={response.status_code}, content_type={response.headers.get('content-type', 'unknown')}, size={len(response.content)} bytes")
            print(f"📸 {api['name']}: status={response.status_code}, size={len(response.content)} bytes")
            
            if response.status_code == 200:
                # Google PageSpeed API の特別処理
                if api.get('extract_screenshot'):
                    try:
                        import base64
                        data = response.json()
                        screenshot_data = data.get('lighthouseResult', {}).get('audits', {}).get('final-screenshot', {}).get('details', {}).get('data', '')
                        
                        if screenshot_data and screenshot_data.startswith('data:image'):
                            # data:image/jpeg;base64,... から base64 部分を抽出
                            base64_str = screenshot_data.split(',')[1]
                            image_data = base64.b64decode(base64_str)
                            logger.info(f"✅ Extracted screenshot from PageSpeed: {len(image_data)} bytes")
                            print(f"✅ PageSpeed screenshot extracted: {len(image_data)} bytes")
                        else:
                            logger.warning(f"❌ PageSpeed: no screenshot data found")
                            print(f"❌ PageSpeed: no screenshot data")
                            continue
                    except Exception as extract_error:
                        logger.warning(f"❌ PageSpeed extraction error: {extract_error}")
                        print(f"❌ PageSpeed extraction failed")
                        continue
                else:
                    image_data = response.content
                
                # 画像が有効かチェック（最低10KB以上）
                if len(image_data) > 10000:
                    logger.info(f"✅ Screenshot captured via {api['name']}: {len(image_data)} bytes")
                    print(f"✅ Screenshot success via {api['name']}: {len(image_data)} bytes")
                    return io.BytesIO(image_data)
                else:
                    logger.warning(f"❌ {api['name']} returned small image: {len(image_data)} bytes (probably blank)")
                    print(f"❌ {api['name']}: image too small ({len(image_data)} bytes)")
            else:
                logger.warning(f"❌ {api['name']} failed: status={response.status_code}, response={response.text[:200]}")
                print(f"❌ {api['name']}: status {response.status_code}")
                
        except Exception as e:
            logger.warning(f"❌ {api['name']} error: {e}")
            print(f"❌ {api['name']} error: {str(e)[:80]}")
            import traceback
            logger.warning(traceback.format_exc()[:500])
            continue
    
    # すべてのAPIが失敗
    logger.error(f"❌ All screenshot APIs failed for: {url}")
    print(f"❌ All screenshot APIs failed")
    return None

def capture_firework_element_screenshot(url, width=1200, height=800, timeout=15):
    """
    Playwrightを使用してFirework要素全体のスクリーンショットを撮影（短時間・要素指定）
    
    Args:
        url: スクリーンショット対象のURL
        width: ビューポート幅
        height: ビューポート高さ
        timeout: タイムアウト（秒）- デフォルト15秒
        
    Returns:
        BytesIO object containing PNG image, or None if failed
    """
    print(f"🎯 capture_firework_element_screenshot CALLED: url={url}, timeout={timeout}s")
    
    try:
        from playwright.sync_api import sync_playwright, TimeoutError as PlaywrightTimeoutError
        
        logger.info(f"🎯 Starting Playwright for Firework element screenshot: {url}")
        
        with sync_playwright() as p:
            # Chromiumブラウザを起動（最小限の設定で高速化）
            browser = p.chromium.launch(
                headless=True,
                args=[
                    '--disable-blink-features=AutomationControlled',
                    '--disable-dev-shm-usage',
                    '--no-sandbox'
                ]
            )
            
            context = browser.new_context(
                viewport={'width': width, 'height': height},
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            )
            
            page = context.new_page()
            
            try:
                # ページを読み込む（短時間タイムアウト）
                logger.info(f"Loading page with {timeout}s timeout...")
                page.goto(url, wait_until='domcontentloaded', timeout=timeout * 1000)
                page.wait_for_timeout(2000)  # 2秒待機してJavaScriptを実行
                
                # Firework要素を探す
                selectors = [
                    'fw-embed-feed',
                    'fw-storyblock',
                    'fw-video-player',
                    '[class*="firework"]',
                    '[id*="firework"]'
                ]
                
                firework_element = None
                for selector in selectors:
                    try:
                        element = page.query_selector(selector)
                        if element:
                            firework_element = element
                            logger.info(f"✅ Found Firework element: {selector}")
                            print(f"✅ Found Firework element: {selector}")
                            break
                    except Exception:
                        continue
                
                if firework_element:
                    # 要素までスクロール（短時間タイムアウト）
                    try:
                        firework_element.scroll_into_view_if_needed(timeout=5000)
                        page.wait_for_timeout(1000)
                    except Exception as scroll_error:
                        logger.warning(f"Scroll failed, continuing: {scroll_error}")
                    
                    # 要素のスクリーンショットを撮る
                    screenshot_bytes = firework_element.screenshot(type='png', timeout=5000)
                    
                    browser.close()
                    
                    logger.info(f"✅ Firework element screenshot captured: {len(screenshot_bytes)} bytes")
                    print(f"✅ Firework element screenshot: {len(screenshot_bytes)} bytes")
                    return io.BytesIO(screenshot_bytes)
                else:
                    logger.warning("❌ No Firework elements found on page")
                    print("❌ No Firework elements found")
                    browser.close()
                    return None
                    
            except PlaywrightTimeoutError as timeout_error:
                logger.warning(f"⏱️ Playwright timeout: {timeout_error}")
                print(f"⏱️ Playwright timeout")
                browser.close()
                return None
                
            except Exception as page_error:
                logger.warning(f"❌ Page error: {page_error}")
                print(f"❌ Page error: {str(page_error)[:80]}")
                browser.close()
                return None
                
    except Exception as e:
        logger.error(f"❌ Playwright error: {e}")
        print(f"❌ Playwright error: {str(e)[:80]}")
        return None

def capture_screenshot_with_playwright(url, width=1200, height=800, firework_format=None):
    """❌ この関数は使用禁止です（タイムアウトのため）
    
    代わりに capture_screenshot_with_api() を使用してください。
    """
    logger.error("❌ capture_screenshot_with_playwright is DEPRECATED - use capture_screenshot_with_api instead")
    print("❌ Playwright function called - this should not happen!")
    return None
    try:
        from playwright.sync_api import sync_playwright, TimeoutError as PlaywrightTimeoutError
        
        logger.info(f"Capturing screenshot for: {url}, format: {firework_format}")
        
        with sync_playwright() as p:
            # Chromiumブラウザを起動（ヘッドレスモード + アクセス強化）
            browser = p.chromium.launch(
                headless=True,
                args=[
                    '--disable-blink-features=AutomationControlled',  # ボット検出回避
                    '--disable-dev-shm-usage',  # メモリ不足対策
                    '--no-sandbox',  # サンドボックス無効化
                    '--disable-setuid-sandbox',  # セキュリティサンドボックス無効化
                    '--disable-web-security',  # ⚡ アクセス強化: Web Security無効化
                    '--disable-features=IsolateOrigins,site-per-process',  # ⚡ アクセス強化: CORS回避
                    '--disable-site-isolation-trials',  # ⚡ アクセス強化: サイト分離無効化
                ]
            )
            
            context = browser.new_context(
                viewport={'width': width, 'height': height},
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                locale='ja-JP',
                timezone_id='Asia/Tokyo',
                ignore_https_errors=True,  # SSL証明書エラーを無視
                bypass_csp=True,  # ⚡ アクセス強化: CSP（Content Security Policy）をバイパス
                java_script_enabled=True,  # JavaScript有効化
            )
            
            # ⚡ アクセス強化: リクエストヘッダーにカスタムヘッダーを追加
            context.set_extra_http_headers({
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8',
                'Accept-Language': 'ja,en-US;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate, br',
                'Connection': 'keep-alive',
                'Upgrade-Insecure-Requests': '1',
                'Sec-Fetch-Dest': 'document',
                'Sec-Fetch-Mode': 'navigate',
                'Sec-Fetch-Site': 'none',
            })
            
            page = context.new_page()
            
            # ⚡ アクセス強化: ボット検出回避のためのJavaScript注入
            page.add_init_script("""
                // Webdriverプロパティを隠蔽
                Object.defineProperty(navigator, 'webdriver', {
                    get: () => undefined
                });
                
                // Chrome特有のプロパティを追加
                window.chrome = {
                    runtime: {}
                };
                
                // Permissions APIをモック
                const originalQuery = window.navigator.permissions.query;
                window.navigator.permissions.query = (parameters) => (
                    parameters.name === 'notifications' ?
                        Promise.resolve({ state: Notification.permission }) :
                        originalQuery(parameters)
                );
                
                // Plugin配列を追加
                Object.defineProperty(navigator, 'plugins', {
                    get: () => [1, 2, 3, 4, 5]
                });
                
                // 言語設定
                Object.defineProperty(navigator, 'languages', {
                    get: () => ['ja-JP', 'ja', 'en-US', 'en']
                });
            """)
            
            # 複数の戦略でリトライ（domcontentloadedを優先）
            strategies = [
                {'wait_until': 'domcontentloaded', 'timeout': 20000},  # 戦略1: DOMContentLoaded（最も速い）
                {'wait_until': 'load', 'timeout': 20000},  # 戦略2: load
                {'wait_until': 'networkidle', 'timeout': 25000},  # 戦略3: networkidle（最も厳格だが遅い）
            ]
            
            screenshot_bytes = None
            last_error = None
            
            for i, strategy in enumerate(strategies, 1):
                try:
                    logger.info(f"Screenshot attempt {i}/3: wait_until={strategy['wait_until']}, timeout={strategy['timeout']}ms")
                    
                    # URLにアクセス
                    page.goto(url, **strategy)
                    
                    # 人間のように振る舞う：マウス移動とスクロール
                    try:
                        # ランダムなマウス移動
                        page.mouse.move(100, 100)
                        page.wait_for_timeout(300)
                        page.mouse.move(500, 300)
                        page.wait_for_timeout(300)
                        
                        # スクロールダウン（ゆっくり）
                        for _ in range(3):
                            page.evaluate('window.scrollBy(0, 300)')
                            page.wait_for_timeout(200)
                        
                        # 最上部に戻る
                        page.evaluate('window.scrollTo(0, 0)')
                        page.wait_for_timeout(500)
                        
                        logger.info("✅ Human-like behavior simulation completed")
                    except Exception as behavior_error:
                        logger.warning(f"Human behavior simulation failed: {behavior_error}")
                    
                    # 少し待機してページを安定させる
                    page.wait_for_timeout(2000)  # 2秒に短縮
                    
                    # ⚡ 強力なポップアップ・モーダルクロージング
                    try:
                        logger.info("Aggressively closing all popups and modals...")
                        
                        # JavaScriptでポップアップを強制的に削除（より積極的）
                        page.evaluate("""
                            () => {
                                // 1. 固定位置の要素を削除（ポップアップやモーダル）
                                const fixedElements = Array.from(document.querySelectorAll('*')).filter(el => {
                                    const style = window.getComputedStyle(el);
                                    return style.position === 'fixed' || style.position === 'absolute';
                                });
                                fixedElements.forEach(el => {
                                    // Firework要素は保持
                                    const tagName = el.tagName.toLowerCase();
                                    if (!tagName.startsWith('fw-')) {
                                        const zIndex = parseInt(window.getComputedStyle(el).zIndex);
                                        // z-indexが高い要素（ポップアップの可能性）を削除
                                        if (zIndex > 100) {
                                            el.style.display = 'none';
                                            el.remove();
                                        }
                                    }
                                });
                                
                                // 2. z-indexが非常に高い要素を削除
                                const highZIndexElements = Array.from(document.querySelectorAll('*')).filter(el => {
                                    const zIndex = parseInt(window.getComputedStyle(el).zIndex);
                                    return zIndex > 1000;
                                });
                                highZIndexElements.forEach(el => {
                                    if (!el.tagName.toLowerCase().startsWith('fw-')) {
                                        el.style.display = 'none';
                                        el.remove();
                                    }
                                });
                                
                                // 3. 一般的なモーダル・オーバーレイクラスを削除
                                const modalSelectors = [
                                    '.modal', '.popup', '.overlay', '.dialog',
                                    '.modal-overlay', '.popup-overlay', '.modal-backdrop',
                                    '[class*="overlay"]', '[id*="overlay"]',
                                    '[class*="modal"]', '[id*="modal"]',
                                    '[class*="popup"]', '[id*="popup"]',
                                    '[class*="dialog"]', '[id*="dialog"]'
                                ];
                                modalSelectors.forEach(selector => {
                                    try {
                                        const elements = document.querySelectorAll(selector);
                                        elements.forEach(el => {
                                            if (!el.tagName.toLowerCase().startsWith('fw-')) {
                                                el.style.display = 'none';
                                                el.remove();
                                            }
                                        });
                                    } catch (e) {}
                                });
                                
                                // 4. WorldShopping関連を削除
                                const wsElements = document.querySelectorAll('[class*="ws-"], [id*="ws-"], [class*="worldshopping"], [id*="worldshopping"]');
                                wsElements.forEach(el => {
                                    el.style.display = 'none';
                                    el.remove();
                                });
                                
                                // 5. iframeを削除（Firework以外）
                                const iframes = document.querySelectorAll('iframe');
                                iframes.forEach(iframe => {
                                    const src = iframe.src || '';
                                    if (!src.includes('firework') && !src.includes('fw-')) {
                                        iframe.style.display = 'none';
                                        iframe.remove();
                                    }
                                });
                                
                                // 6. bodyのスクロールを有効化し、overflow: hiddenを解除
                                document.body.style.overflow = 'auto !important';
                                document.body.style.position = 'static';
                                document.documentElement.style.overflow = 'auto !important';
                                document.documentElement.style.position = 'static';
                                
                                // 7. 半透明の背景要素を削除（opacity < 1 かつ大きい要素）
                                Array.from(document.querySelectorAll('*')).forEach(el => {
                                    const style = window.getComputedStyle(el);
                                    const opacity = parseFloat(style.opacity);
                                    const width = el.offsetWidth;
                                    const height = el.offsetHeight;
                                    // 画面サイズより大きく、半透明の要素はオーバーレイの可能性
                                    if (opacity < 1 && opacity > 0 && width > window.innerWidth * 0.8 && height > window.innerHeight * 0.8) {
                                        if (!el.tagName.toLowerCase().startsWith('fw-')) {
                                            el.style.display = 'none';
                                            el.remove();
                                        }
                                    }
                                });
                                
                                // 🆕 8. テキストベースのフィルタリング: 特定のテキストを含む要素を削除
                                const popupTexts = [
                                    '海外にお住まいのお客様へ',
                                    '海外にお住まいのお客様',
                                    'お住まいのお客様へ',
                                    'Cookie',
                                    'クッキー',
                                    '個人情報保護方針',
                                    'プライバシーポリシー',
                                    '同意する',
                                    'Accept',
                                    '閉じる',
                                    'Close',
                                    'WorldShopping'
                                ];
                                
                                Array.from(document.querySelectorAll('*')).forEach(el => {
                                    if (el.tagName.toLowerCase().startsWith('fw-')) return;
                                    
                                    const text = el.textContent || '';
                                    const innerText = el.innerText || '';
                                    
                                    for (const popupText of popupTexts) {
                                        if (text.includes(popupText) || innerText.includes(popupText)) {
                                            const style = window.getComputedStyle(el);
                                            if (style.position === 'fixed' || 
                                                style.position === 'absolute' || 
                                                parseInt(style.zIndex) > 50) {
                                                el.style.display = 'none';
                                                el.remove();
                                                console.log('Removed popup with text:', popupText);
                                                break;
                                            }
                                        }
                                    }
                                });
                            }
                        """)
                        
                        logger.info("Popups removed via JavaScript (1st pass)")
                        # ポップアップ削除後、少し待機してDOMを安定させる
                        page.wait_for_timeout(2000)
                        
                        # 🆕 2回目のポップアップ削除（遅延表示されるポップアップ対策）
                        logger.info("Running 2nd popup removal pass for delayed popups...")
                        page.evaluate("""
                            () => {
                                // テキストベースのフィルタリング（2回目）
                                const popupTexts = [
                                    '海外にお住まいのお客様へ',
                                    '海外にお住まいのお客様',
                                    'お住まいのお客様へ',
                                    'Cookie', 'クッキー',
                                    '個人情報保護方針',
                                    'プライバシーポリシー',
                                    '同意する', 'Accept',
                                    '閉じる', 'Close',
                                    'WorldShopping'
                                ];
                                
                                Array.from(document.querySelectorAll('*')).forEach(el => {
                                    if (el.tagName.toLowerCase().startsWith('fw-')) return;
                                    
                                    const text = el.textContent || '';
                                    const innerText = el.innerText || '';
                                    
                                    for (const popupText of popupTexts) {
                                        if (text.includes(popupText) || innerText.includes(popupText)) {
                                            const style = window.getComputedStyle(el);
                                            if (style.position === 'fixed' || 
                                                style.position === 'absolute' || 
                                                parseInt(style.zIndex) > 50) {
                                                el.style.display = 'none';
                                                el.remove();
                                                console.log('Removed delayed popup with text:', popupText);
                                                break;
                                            }
                                        }
                                    }
                                });
                                
                                // 固定位置の高z-index要素も再削除
                                Array.from(document.querySelectorAll('*')).forEach(el => {
                                    if (el.tagName.toLowerCase().startsWith('fw-')) return;
                                    const style = window.getComputedStyle(el);
                                    const zIndex = parseInt(style.zIndex);
                                    if ((style.position === 'fixed' || style.position === 'absolute') && zIndex > 500) {
                                        el.style.display = 'none';
                                        el.remove();
                                    }
                                });
                            }
                        """)
                        logger.info("✅ 2nd popup removal pass complete")
                        page.wait_for_timeout(1000)
                        
                    except Exception as popup_error:
                        logger.warning(f"JavaScript popup removal failed: {popup_error}")
                    
                    # Fireworkフォーマットが指定されている場合、そのフォーマットの要素を探す
                    if firework_format and firework_format != 'Unknown':
                        try:
                            logger.info(f"Looking for specific Firework format: {firework_format}")
                            
                            # フォーマット名に基づいて適切な要素を探す
                            format_to_selector_map = {
                                'Horizontal Carousel': {'selector': 'fw-embed-feed', 'attribute': 'style', 'value': 'thumbnail'},
                                'Dynamic Carousel': {'selector': 'fw-embed-feed', 'attribute': 'thumbnail_style', 'value': 'dynamic'},
                                'Grid': {'selector': 'fw-embed-feed', 'attribute': 'mode', 'value': 'grid'},
                                'Carousel': {'selector': 'fw-embed-feed', 'attribute': 'mode', 'value': 'row'},
                                'Story Block': {'selector': 'fw-embed-feed', 'exclude_attrs': ['mode', 'style', 'thumbnail_style']},
                                'Circle Stories': {'selector': 'fw-stories', 'attribute': 'thumbnail_shape', 'value': 'circle'},
                                'Vertical Stories': {'selector': 'fw-stories', 'attribute': 'thumbnail_shape', 'value': 'rectangle'},
                                'Floating Player': {'selector': 'fw-storyblock', 'attribute': 'mode', 'value': 'pinned'},
                                'Horizontal Player': {'selector': 'fw-player'},
                                'Hero Unit': {'selector': 'fw-herounit'},
                                'Player Deck': {'selector': 'fw-player-deck'},
                            }
                            
                            element_found = None
                            
                            if firework_format in format_to_selector_map:
                                format_config = format_to_selector_map[firework_format]
                                selector = format_config['selector']
                                
                                logger.info(f"Searching for <{selector}> elements...")
                                
                                # JavaScriptでFirework要素を直接探す（より確実）
                                # 引数を1つのオブジェクトにまとめる（Playwright制限対応）
                                search_params = {
                                    'selector': selector,
                                    'attrName': format_config.get('attribute'),
                                    'attrValue': format_config.get('value'),
                                    'excludeAttrs': format_config.get('exclude_attrs', [])
                                }
                                
                                matching_elements = page.evaluate("""
                                    (params) => {
                                        const elements = Array.from(document.querySelectorAll(params.selector));
                                        const results = [];
                                        
                                        elements.forEach((el, index) => {
                                            const isVisible = el.offsetWidth > 0 && el.offsetHeight > 0;
                                            if (!isVisible) return;
                                            
                                            const outerHTML = el.outerHTML;
                                            let matches = false;
                                            
                                            if (params.attrName) {
                                                if (params.attrName === 'style') {
                                                    matches = outerHTML.includes(params.attrValue);
                                                } else {
                                                    const attrVal = el.getAttribute(params.attrName);
                                                    matches = attrVal === params.attrValue;
                                                }
                                            } else if (params.excludeAttrs && params.excludeAttrs.length > 0) {
                                                matches = !params.excludeAttrs.some(attr => el.hasAttribute(attr));
                                            } else {
                                                matches = true;
                                            }
                                            
                                            if (matches) {
                                                const rect = el.getBoundingClientRect();
                                                results.push({
                                                    index: index,
                                                    top: rect.top,
                                                    left: rect.left,
                                                    width: rect.width,
                                                    height: rect.height,
                                                    outerHTML: outerHTML.substring(0, 200)
                                                });
                                            }
                                        });
                                        
                                        return results;
                                    }
                                """, search_params)
                                
                                logger.info(f"JavaScript found {len(matching_elements)} matching elements")
                                
                                # 要素が見つからない場合、追加で待機してリトライ
                                if len(matching_elements) == 0:
                                    logger.warning("No elements found on first attempt, waiting 8 more seconds...")
                                    page.wait_for_timeout(8000)  # 5秒→8秒に延長
                                    
                                    # 再度検索
                                    matching_elements = page.evaluate("""
                                        (params) => {
                                            const elements = Array.from(document.querySelectorAll(params.selector));
                                            const results = [];
                                            
                                            elements.forEach((el, index) => {
                                                const isVisible = el.offsetWidth > 0 && el.offsetHeight > 0;
                                                if (!isVisible) return;
                                                
                                                const outerHTML = el.outerHTML;
                                                let matches = false;
                                                
                                                if (params.attrName) {
                                                    if (params.attrName === 'style') {
                                                        matches = outerHTML.includes(params.attrValue);
                                                    } else {
                                                        const attrVal = el.getAttribute(params.attrName);
                                                        matches = attrVal === params.attrValue;
                                                    }
                                                } else if (params.excludeAttrs && params.excludeAttrs.length > 0) {
                                                    matches = !params.excludeAttrs.some(attr => el.hasAttribute(attr));
                                                } else {
                                                    matches = true;
                                                }
                                                
                                                if (matches) {
                                                    const rect = el.getBoundingClientRect();
                                                    results.push({
                                                        index: index,
                                                        top: rect.top,
                                                        left: rect.left,
                                                        width: rect.width,
                                                        height: rect.height,
                                                        outerHTML: outerHTML.substring(0, 200)
                                                    });
                                                }
                                            });
                                            
                                            return results;
                                        }
                                    """, search_params)
                                    logger.info(f"After retry: JavaScript found {len(matching_elements)} matching elements")
                                
                                if len(matching_elements) > 0:
                                    # 最初のマッチする要素を使用
                                    target_info = matching_elements[0]
                                    logger.info(f"Target element: index={target_info['index']}, size={target_info['width']}x{target_info['height']}")
                                    logger.info(f"HTML: {target_info['outerHTML']}")
                                    
                                    # 要素を取得
                                    all_elements = page.locator(selector).all()
                                    if target_info['index'] < len(all_elements):
                                        element_found = all_elements[target_info['index']]
                                        
                                        # 🔥 NEW APPROACH: JavaScriptで要素を確実にビューポートの中央にスクロール
                                        logger.info("Scrolling element to center of viewport using JavaScript...")
                                        scroll_result = element_found.evaluate("""
                                            el => {
                                                // 要素をビューポートの中央にスクロール
                                                el.scrollIntoView({
                                                    behavior: 'auto',  // smooth scrollは使わない（完了を待てない）
                                                    block: 'center',   // 縦方向中央
                                                    inline: 'center'   // 横方向中央
                                                });
                                                
                                                // スクロール後の位置を返す
                                                const rect = el.getBoundingClientRect();
                                                const scrollY = window.pageYOffset || document.documentElement.scrollTop;
                                                
                                                return {
                                                    scrolledTo: scrollY,
                                                    viewportY: rect.top,
                                                    viewportX: rect.left,
                                                    viewportBottom: rect.bottom,
                                                    width: rect.width,
                                                    height: rect.height
                                                };
                                            }
                                        """)
                                        
                                        logger.info(f"Scroll result: viewportY={scroll_result['viewportY']}, height={scroll_result['height']}, scrolledTo={scroll_result['scrolledTo']}")
                                        
                                        # スクロールアニメーション完了を待つ
                                        page.wait_for_timeout(2000)
                                        
                                        # 🎥 Floating Playerの場合、動画がロードされるまで待つ
                                        if firework_format == 'Floating Player':
                                            logger.info("🎥 Waiting for Floating Player video to load...")
                                            try:
                                                # fw-storyblock内のvideo/iframe要素が現れるまで待機
                                                page.wait_for_selector('fw-storyblock video, fw-storyblock iframe', timeout=10000)
                                                logger.info("✅ Video element detected")
                                                # 動画の初期化を待つ（追加で2-3秒）
                                                page.wait_for_timeout(3000)
                                                logger.info("✅ Video loading wait complete")
                                            except Exception as video_wait_error:
                                                logger.warning(f"⚠️ Video wait timeout or error: {video_wait_error}")
                                                # タイムアウトしても続行（動画がない場合もある）
                                                pass
                                        
                                        # 再度ポップアップをJavaScriptで削除（より積極的 + テキストベースフィルタリング）
                                        try:
                                            page.evaluate("""
                                                () => {
                                                    // 固定位置の要素を全て削除
                                                    const fixedElements = Array.from(document.querySelectorAll('*')).filter(el => {
                                                        const style = window.getComputedStyle(el);
                                                        return (style.position === 'fixed' || style.position === 'absolute') && 
                                                               parseInt(style.zIndex) > 100;
                                                    });
                                                    fixedElements.forEach(el => {
                                                        if (!el.tagName.toLowerCase().startsWith('fw-')) {
                                                            el.style.display = 'none';
                                                            el.remove();
                                                        }
                                                    });
                                                    
                                                    // オーバーレイ系を全て削除
                                                    const overlaySelectors = ['.modal', '.popup', '.overlay', '.dialog', '[class*="overlay"]', '[class*="modal"]'];
                                                    overlaySelectors.forEach(selector => {
                                                        try {
                                                            document.querySelectorAll(selector).forEach(el => {
                                                                if (!el.tagName.toLowerCase().startsWith('fw-')) {
                                                                    el.style.display = 'none';
                                                                    el.remove();
                                                                }
                                                            });
                                                        } catch(e) {}
                                                    });
                                                    
                                                    // 🆕 テキストベースのフィルタリング: 特定のテキストを含む要素を削除
                                                    const popupTexts = [
                                                        '海外にお住まいのお客様へ',
                                                        '海外にお住まいのお客様',
                                                        'お住まいのお客様へ',
                                                        'Cookie',
                                                        'クッキー',
                                                        '個人情報保護方針',
                                                        'プライバシーポリシー',
                                                        '同意する',
                                                        'Accept',
                                                        '閉じる',
                                                        'Close'
                                                    ];
                                                    
                                                    // すべての要素をチェック
                                                    Array.from(document.querySelectorAll('*')).forEach(el => {
                                                        // Firework要素はスキップ
                                                        if (el.tagName.toLowerCase().startsWith('fw-')) return;
                                                        
                                                        const text = el.textContent || '';
                                                        const innerText = el.innerText || '';
                                                        
                                                        // ポップアップテキストが含まれているかチェック
                                                        for (const popupText of popupTexts) {
                                                            if (text.includes(popupText) || innerText.includes(popupText)) {
                                                                // この要素またはその親要素を削除
                                                                const style = window.getComputedStyle(el);
                                                                // 固定位置または高いz-indexを持つ場合
                                                                if (style.position === 'fixed' || 
                                                                    style.position === 'absolute' || 
                                                                    parseInt(style.zIndex) > 50) {
                                                                    el.style.display = 'none';
                                                                    el.remove();
                                                                    console.log('Removed popup with text:', popupText);
                                                                    break;
                                                                }
                                                            }
                                                        }
                                                    });
                                                }
                                            """)
                                            logger.info("✅ Popups removed after scroll (with text-based filtering)")
                                        except Exception as e:
                                            logger.warning(f"Popup removal after scroll failed: {e}")
                                        page.wait_for_timeout(500)
                                        
                                        # 🔥 要素を含む周辺コンテキストをスクリーンショット
                                        try:
                                            # 再度要素の位置情報を取得（ポップアップ削除後）
                                            element_info = element_found.evaluate("""
                                                el => {
                                                    const rect = el.getBoundingClientRect();
                                                    const scrollY = window.pageYOffset || document.documentElement.scrollTop;
                                                    const viewportHeight = window.innerHeight;
                                                    
                                                    return {
                                                        viewportY: rect.top,
                                                        viewportX: rect.left,
                                                        width: rect.width,
                                                        height: rect.height,
                                                        scrollY: scrollY,
                                                        viewportHeight: viewportHeight,
                                                        visible: rect.width > 0 && rect.height > 0,
                                                        inViewport: rect.top >= 0 && rect.bottom <= viewportHeight,
                                                        bottom: rect.bottom
                                                    };
                                                }
                                            """)
                                            
                                            logger.info(f"Element position after popup removal: viewportY={element_info['viewportY']}, height={element_info['height']}, visible={element_info['visible']}, inViewport={element_info['inViewport']}, scrollY={element_info['scrollY']}")
                                            
                                            if element_info['visible'] and element_info['height'] > 0:
                                                # ビューポート座標を使用してスクリーンショット
                                                viewport_y = element_info['viewportY']
                                                element_height = element_info['height']
                                                viewport_height = element_info['viewportHeight']
                                                
                                                # パディング（上下に余白を追加）
                                                padding_top = 200  # 上部の余白を増やす
                                                padding_bottom = 200  # 下部の余白を増やす
                                                
                                                # clip座標を計算（ビューポート座標系）
                                                # 要素の上部 - padding_top が開始位置
                                                clip_y = max(0, viewport_y - padding_top)
                                                
                                                # 終了位置は要素の下部 + padding_bottom、ただしビューポート高さを超えない
                                                clip_end_y = min(viewport_y + element_height + padding_bottom, viewport_height)
                                                clip_height = clip_end_y - clip_y
                                                
                                                logger.info(f"Clip calculation: clip_y={clip_y}, clip_height={clip_height}, viewport_height={viewport_height}")
                                                
                                                # clip_heightが有効範囲内にあることを確認
                                                if clip_height > 50 and clip_height <= viewport_height:
                                                    logger.info(f"📸 Capturing screenshot with clip: x=0, y={clip_y}, w={width}, h={clip_height}")
                                                    
                                                    screenshot_bytes = page.screenshot(
                                                        type='png',
                                                        full_page=False,  # ビューポートのみをキャプチャ
                                                        clip={
                                                            'x': 0,
                                                            'y': clip_y,
                                                            'width': width,
                                                            'height': clip_height
                                                        }
                                                    )
                                                    logger.info(f"✅ Screenshot captured successfully: {len(screenshot_bytes)} bytes")
                                                else:
                                                    # clip_heightが無効な場合、ビューポート全体をスクリーンショット
                                                    logger.warning(f"⚠️ Invalid clip_height: {clip_height}, capturing full viewport instead")
                                                    screenshot_bytes = page.screenshot(full_page=False, type='png')
                                            else:
                                                # 要素が見えない場合
                                                logger.error(f"❌ Element not visible: height={element_info['height']}, visible={element_info['visible']}")
                                                logger.warning("Falling back to full viewport screenshot")
                                                screenshot_bytes = page.screenshot(full_page=False, type='png')
                                                    
                                        except Exception as screenshot_error:
                                            logger.error(f"Screenshot error: {screenshot_error}")
                                            logger.error(traceback.format_exc())
                                            # 最終フォールバック
                                            logger.warning("Exception occurred, capturing full viewport as fallback")
                                            screenshot_bytes = page.screenshot(full_page=False, type='png')
                                            logger.info(f"Final fallback screenshot: {len(screenshot_bytes)} bytes")
                                else:
                                    logger.error(f"❌ No matching elements found for format '{firework_format}' after retry")
                                    logger.info("Setting screenshot_bytes to None to trigger HTML file saving")
                                    screenshot_bytes = None  # HTMLファイル保存に進む
                            else:
                                logger.warning(f"Unknown format: {firework_format}, will try HTML file saving")
                                screenshot_bytes = None  # HTMLファイル保存に進む
                        
                        except Exception as fw_error:
                            logger.error(f"Firework element screenshot failed: {fw_error}")
                            logger.error(traceback.format_exc())
                            # フォールバック: HTMLファイル保存に進む
                            logger.info("Setting screenshot_bytes to None due to exception")
                            screenshot_bytes = None
                    else:
                        # フォーマットが指定されていない場合、通常のスクリーンショット
                        screenshot_bytes = page.screenshot(full_page=False, type='png')
                    
                    # screenshot_bytesがNoneでないことを確認してからlenを呼ぶ
                    if screenshot_bytes is not None:
                        logger.info(f"Screenshot captured successfully with strategy {i}: {len(screenshot_bytes)} bytes")
                        break  # 成功したらループを抜ける
                    else:
                        logger.warning(f"Screenshot is None, will try HTML fallback")
                        continue  # 次の戦略は試さず、HTMLフォールバックに進む
                    
                except PlaywrightTimeoutError as timeout_error:
                    last_error = timeout_error
                    logger.warning(f"Strategy {i} timed out: {timeout_error}")
                    continue  # 次の戦略を試す
                    
                except Exception as strategy_error:
                    last_error = strategy_error
                    logger.warning(f"Strategy {i} failed: {strategy_error}")
                    continue
            
            browser.close()
            
            if screenshot_bytes:
                return io.BytesIO(screenshot_bytes)
            else:
                logger.error(f"All screenshot strategies failed for {url}")
                if last_error:
                    logger.error(f"Last error: {last_error}")
                return None
            
    except Exception as e:
        logger.error(f"Playwright screenshot failed: {e}")
        logger.error(traceback.format_exc())
        return None

def capture_firework_video_thumbnail(url, width=400, height=300):
    """
    Firework動画の個別サムネイルを取得（HTMLからFirework要素を検出）
    
    Args:
        url: ページURL
        width: サムネイル幅
        height: サムネイル高さ
        
    Returns:
        BytesIO object containing PNG image, or None if failed
    """
    print(f"🎬 capture_firework_video_thumbnail CALLED: url={url}")
    logger.info(f"🎬 Attempting to extract Firework video thumbnail from HTML")
    
    try:
        from bs4 import BeautifulSoup
        import requests
        
        # HTMLを取得
        response = requests.get(url, timeout=10, headers={
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        })
        
        if response.status_code != 200:
            logger.warning(f"Failed to fetch HTML: status={response.status_code}")
            return None
        
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Firework要素を検索（複数のタグパターンに対応）
        firework_selectors = [
            'fw-storyblock',
            'fw-video-player', 
            'fw-embed-feed',
            '[class*="firework"]',
            '[id*="firework"]'
        ]
        
        firework_element = None
        for selector in firework_selectors:
            elements = soup.select(selector)
            if elements:
                firework_element = elements[0]  # 最初の要素を取得
                logger.info(f"✅ Found Firework element: {selector}")
                print(f"✅ Found Firework element: {selector}")
                break
        
        if not firework_element:
            logger.warning("❌ No Firework elements found in HTML")
            print("❌ No Firework elements found")
            # フォールバック: ページ全体のスクリーンショット
            return capture_screenshot_with_api(url, width=width, height=height)
        
        # サムネイル画像URLを抽出（複数のパターンを試す）
        thumbnail_url = None
        
        # パターン1: poster属性
        if firework_element.get('poster'):
            thumbnail_url = firework_element.get('poster')
            logger.info(f"Found thumbnail via poster: {thumbnail_url}")
        
        # パターン2: data-video-url属性
        elif firework_element.get('data-video-url'):
            thumbnail_url = firework_element.get('data-video-url')
            logger.info(f"Found thumbnail via data-video-url: {thumbnail_url}")
        
        # パターン3: img要素を探す
        elif firework_element.find('img'):
            img = firework_element.find('img')
            thumbnail_url = img.get('src') or img.get('data-src')
            logger.info(f"Found thumbnail via img tag: {thumbnail_url}")
        
        # サムネイルURLが見つかった場合、画像を直接ダウンロード
        if thumbnail_url:
            # 相対URLを絶対URLに変換
            from urllib.parse import urljoin
            thumbnail_url = urljoin(url, thumbnail_url)
            
            logger.info(f"📥 Downloading Firework thumbnail: {thumbnail_url}")
            print(f"📥 Downloading thumbnail: {thumbnail_url[:80]}...")
            
            img_response = requests.get(thumbnail_url, timeout=10, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            })
            
            if img_response.status_code == 200 and len(img_response.content) > 1000:
                logger.info(f"✅ Firework thumbnail downloaded: {len(img_response.content)} bytes")
                print(f"✅ Thumbnail downloaded: {len(img_response.content)} bytes")
                return io.BytesIO(img_response.content)
            else:
                logger.warning(f"❌ Failed to download thumbnail: status={img_response.status_code}")
        
        # すべて失敗した場合、ページ全体のスクリーンショットにフォールバック
        logger.warning("⚠️ No valid thumbnail found, falling back to full page screenshot")
        print("⚠️ Falling back to full page screenshot")
        return capture_screenshot_with_api(url, width=width, height=height)
        
    except Exception as e:
        logger.error(f"Firework thumbnail extraction error: {e}")
        print(f"❌ Thumbnail extraction error: {str(e)[:80]}")
        # エラー時はページ全体のスクリーンショットにフォールバック
        return capture_screenshot_with_api(url, width=width, height=height)

def capture_firework_video_thumbnail_with_playwright(url, width=400, height=300):
    """❌ この関数は使用禁止です（タイムアウトのため）
    
    代わりに capture_firework_video_thumbnail() を使用してください（内部で外部APIを呼び出します）。
    """
    logger.error("❌ capture_firework_video_thumbnail_with_playwright is DEPRECATED - use capture_firework_video_thumbnail instead")
    print("❌ Playwright video thumbnail function called - this should not happen!")
    return None
    try:
        from playwright.sync_api import sync_playwright
        
        logger.info(f"🎬 Starting Firework video thumbnail capture for: {url}")
        
        with sync_playwright() as p:
            browser = p.chromium.launch(
                headless=True,
                args=[
                    '--disable-blink-features=AutomationControlled',
                    '--disable-features=IsolateOrigins,site-per-process',
                    '--disable-web-security',
                    '--ignore-certificate-errors'
                ]
            )
            
            context = browser.new_context(
                viewport={'width': 1280, 'height': 720},
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                locale='ja-JP',
                timezone_id='Asia/Tokyo'
            )
            
            page = context.new_page()
            
            try:
                # ページを読み込む（複数の戦略を試す）
                logger.info(f"Loading page: {url}")
                
                # 戦略1: domcontentloaded（最速）
                try:
                    logger.info("Strategy 1: Trying domcontentloaded...")
                    page.goto(url, wait_until='domcontentloaded', timeout=10000)
                    page.wait_for_timeout(2000)  # 2秒待機してJavaScriptを実行させる
                    logger.info("✅ Strategy 1 succeeded")
                except Exception as e1:
                    logger.warning(f"Strategy 1 failed: {e1}")
                    
                    # 戦略2: load
                    try:
                        logger.info("Strategy 2: Trying load...")
                        page.goto(url, wait_until='load', timeout=15000)
                        page.wait_for_timeout(2000)
                        logger.info("✅ Strategy 2 succeeded")
                    except Exception as e2:
                        logger.warning(f"Strategy 2 failed: {e2}")
                        
                        # 戦略3: commit（最も緩い）
                        logger.info("Strategy 3: Trying commit (most lenient)...")
                        page.goto(url, wait_until='commit', timeout=10000)
                        page.wait_for_timeout(3000)
                        logger.info("✅ Strategy 3 succeeded")
                
                # Firework要素を探す
                # fw-embed-feed, fw-storyblock, fw-video-player などを検索
                selectors = [
                    'fw-embed-feed',
                    'fw-storyblock', 
                    'fw-video-player',
                    '[class*="firework"]',
                    '[id*="firework"]'
                ]
                
                video_element = None
                for selector in selectors:
                    try:
                        elements = page.query_selector_all(selector)
                        if elements and len(elements) > 0:
                            video_element = elements[0]
                            logger.info(f"✅ Found Firework element: {selector}")
                            break
                    except Exception:
                        continue
                
                if video_element:
                    # 要素までスクロール
                    video_element.scroll_into_view_if_needed()
                    page.wait_for_timeout(1000)
                    
                    # 要素のスクリーンショットを撮る
                    screenshot_bytes = video_element.screenshot(type='png')
                    
                    logger.info(f"✅ Firework video thumbnail captured: {len(screenshot_bytes)} bytes")
                    browser.close()
                    
                    # リサイズ
                    from PIL import Image
                    img = Image.open(io.BytesIO(screenshot_bytes))
                    img.thumbnail((width, height), Image.Resampling.LANCZOS)
                    
                    output = io.BytesIO()
                    img.save(output, format='PNG')
                    output.seek(0)
                    
                    return output
                else:
                    logger.warning("⚠️ No Firework video elements found on page")
                    
                    # フォールバック: ページ全体のスクリーンショットを撮って、動画っぽい部分を探す
                    screenshot_bytes = page.screenshot(type='png', full_page=False)
                    logger.info(f"📸 Fallback: captured full viewport screenshot")
                    
                    browser.close()
                    
                    from PIL import Image
                    img = Image.open(io.BytesIO(screenshot_bytes))
                    img.thumbnail((width, height), Image.Resampling.LANCZOS)
                    
                    output = io.BytesIO()
                    img.save(output, format='PNG')
                    output.seek(0)
                    
                    return output
                    
            except Exception as page_error:
                logger.error(f"Page processing error: {page_error}")
                logger.error(traceback.format_exc())
                browser.close()
                return None
                
    except Exception as e:
        logger.error(f"Firework video thumbnail capture failed: {e}")
        logger.error(traceback.format_exc())
        return None

def generate_why_firework(url, html_content, website_description, language='ja', firework_format='Unknown'):
    """OpenAI APIを使用してFirework活用理由を生成（Firework要素とフォーマット情報を含む高度な分析）"""
    try:
        openai_api_key = os.environ.get('OPENAI_API_KEY', '')
        
        if not openai_api_key:
            logger.warning("OPENAI_API_KEY not set")
            fallback = '目的: 動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率' if language == 'ja' else 'Objective: Explain product features and usage through video / Key KPI: Video completion rate'
            return fallback
        
        from bs4 import BeautifulSoup
        
        # HTMLからFirework要素周辺のコンテキストを抽出
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Firework要素を探す
        firework_elements = soup.find_all(lambda tag: tag.name and tag.name.startswith('fw-'))
        firework_context = ""
        
        if firework_elements:
            logger.info(f"Found {len(firework_elements)} Firework elements")
            # Firework要素の周辺テキストを抽出（親要素やsiblingから）
            for fw_elem in firework_elements[:3]:  # 最初の3つのみ
                # 親要素のテキストを取得
                parent = fw_elem.parent
                if parent:
                    parent_text = parent.get_text(separator=' ', strip=True)
                    if parent_text and len(parent_text) > 20:
                        firework_context += parent_text[:200] + " "
                
                # Firework要素の属性情報も取得（channel, playlist情報など）
                attrs = fw_elem.attrs
                if 'channel' in attrs:
                    firework_context += f"[Firework Channel: {attrs['channel']}] "
                if 'playlist' in attrs:
                    firework_context += f"[Playlist: {attrs['playlist']}] "
        
        # HTMLからテキストコンテンツを抽出（一般的なページ内容）
        for script in soup(["script", "style"]):
            script.decompose()
        
        text_content = soup.get_text()
        lines = (line.strip() for line in text_content.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        # テキストが長すぎる場合は切り詰め
        if len(text) > 2000:
            text = text[:2000]
        
        # Fireworkフォーマット情報
        format_description = ""
        if firework_format and firework_format != 'Unknown':
            format_description = f"\n\nFireworkフォーマット: {firework_format}"
        
        # 目的とKPIのパターンリスト
        patterns = [
            "着回し提案で滞在時間を延長 / 主要KPI: 平均滞在時間増加率",
            "素材動画で不安を払拭 / 主要KPI: 離脱率低減",
            "機能実演でCVR向上 / 主要KPI: CVRリフト (非視聴者比)",
            "商品ページの伝達力向上 / 主要KPI: SKUカバー率",
            "ライブ動画を通じて即時購入促進 / 主要KPI: 売上増加",
            "限定品のライブ販売で売上最大化 / 主要KPI: 購入件数・売上",
            "ライブでスタッフのファン化と送客 / 主要KPI: 実店舗来店率",
            "サイマル配信で認知拡大 / 主要KPI: ライブ視聴完了率",
            "ライブ配信でブランド認知向上 / 主要KPI: 検索数増加率",
            "ライブ配信でブランドエンゲージメント向上 / 主要KPI: ライブ中のコメント数増加率",
            "動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率",
            "ライブ配信で新規ユーザーとのエンゲージメントを図る / 主要KPI: 新規視聴者数・新規視聴者率",
            "安心のアフターケアを動画で訴求 / 主要KPI: 申込み率",
            "悩み解決動画でカート追加促進 / 主要KPI: カート追加率",
            "ライブ配信でセット購入促進 / 主要KPI: セット購入率",
            "サイズ感説明動画で障壁低減 / 主要KPI: CVRリフト (非視聴者比)",
            "多様な利用シーンを動画で紹介 / 主要KPI: まとめ買い購入率",
            "商品の操作性を動画で解説 / 主要KPI: CVRリフト (非視聴者比)",
            "設置イメージ動画で決定促進 / 主要KPI: カート追加リフト (非視聴者比)",
            "AIFAQで疑問即時解消 / 主要KPI: チャット満足度",
            "動画を活用し商品ページのコンテンツをリッチ化 / 主要KPI: サイト滞在時間",
            "AIFAQで人的負荷を軽減 / 主要KPI: エスカレーション率低減",
            "AIFAQ分析結果を商品開発へ活用 / 主要KPI: 質問内容（定性）",
            "ライブの双方向性でインサイト獲得 / 主要KPI: コメントの内容（定性評価）",
            "パーソナライズ動画レコメンド強化 / 主要KPI: 平均サイト滞在時間",
            "カテゴリ横断動画でまとめ買い促進 / 主要KPI: セット購入平均点数",
            "動画導入でデジタル体験向上 / 主要KPI: LTV増加率"
        ]
        
        patterns_text = "\n".join([f"{i+1}. {p}" for i, p in enumerate(patterns)])
        
        prompt = f"""以下のウェブサイトの情報とFireworkの動画配置状況を分析し、Fireworkの動画ソリューション活用について最も関連性の高い「目的とKPI」のパターンを1つ選んでください。

ウェブサイトの概要:
{website_description}

Firework動画周辺のコンテキスト:
{firework_context if firework_context else 'Firework要素周辺のコンテンツなし'}{format_description}

ウェブサイトのコンテンツ（一部）:
{text[:1000]}

利用可能な目的とKPIパターン:
{patterns_text}

指示:
1. Firework動画がどのような目的で配置されているか、ページコンテキストとフォーマット情報から推測してください
2. 上記のパターンから最も関連性が高いものを1つ選択してください
3. 選択したパターンをベースに、このウェブサイト固有の状況に合わせてカスタマイズした文章を作成してください
4. 出力は必ず「{'目的: ' if language == 'ja' else 'Objective: '}」で始めてください
5. 出力は80-120文字程度で、「{'目的: ' if language == 'ja' else 'Objective: '}目的内容 / 主要KPI: KPI名」の形式で記述してください
6. {'日本語' if language == 'ja' else '英語'}で出力してください

例:
- アパレルブランドの場合: 「目的: 着回し提案動画で滞在時間を延長し、購入検討を促進 / 主要KPI: 平均滞在時間増加率」
- 家電メーカーの場合: 「目的: 操作性を動画で分かりやすく解説し購入不安を解消 / 主要KPI: CVRリフト (非視聴者比)」

出力（カスタマイズされた1文のみ）:"""

        response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {openai_api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o-mini',
                'messages': [{'role': 'user', 'content': prompt}],
                'max_tokens': 200,
                'temperature': 0.7
            },
            timeout=30
        )
        
        if response.status_code == 200:
            try:
                result = response.json()
                why_firework = result['choices'][0]['message']['content'].strip()
                
                # 「目的: 」または「Objective: 」プレフィックスがない場合は追加
                if language == 'ja':
                    if not why_firework.startswith('目的:'):
                        why_firework = '目的: ' + why_firework
                else:
                    if not why_firework.startswith('Objective:'):
                        why_firework = 'Objective: ' + why_firework
                
                logger.info(f"Why firework generated: {why_firework}")
                return why_firework
            except (ValueError, KeyError) as json_error:
                logger.error(f"Failed to parse Why firework response: {json_error}")
                fallback = '目的: 動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率' if language == 'ja' else 'Objective: Explain product features and usage through video / Key KPI: Video completion rate'
                return fallback
        else:
            logger.error(f"Why firework API error: {response.status_code}")
            logger.error(f"Response: {response.text[:500]}")
            fallback = '目的: 動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率' if language == 'ja' else 'Objective: Explain product features and usage through video / Key KPI: Video completion rate'
            return fallback
            
    except Exception as e:
        logger.error(f"Why firework generation error: {e}")
        logger.error(traceback.format_exc())
        fallback = '目的: 動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率' if language == 'ja' else 'Objective: Explain product features and usage through video / Key KPI: Video completion rate'
        return fallback

def crawl_and_analyze_website(url, language='ja'):
    """requestsとBeautifulSoupを使用してWebサイト情報を取得し、OpenAI APIで分析（高速・安定版）"""
    print(f"🌐🌐🌐 crawl_and_analyze_website CALLED: url={url}, language={language}")
    try:
        from bs4 import BeautifulSoup
        
        fallback = '手動でサイト概要を入力してください' if language == 'ja' else 'Please manually enter website description here'
        
        logger.info(f"🌐 Starting website analysis with requests for: {url}")
        print(f"✅ Starting fast website crawling (no Playwright)...")
        
        # requests でウェブサイトを取得（高速・安定）
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'ja,en-US;q=0.9,en;q=0.8',
                'Accept-Encoding': 'gzip, deflate',
                'Connection': 'keep-alive',
            }
            
            response = requests.get(url, timeout=10, headers=headers, allow_redirects=True)
            response.raise_for_status()
            html_content = response.text
            
            logger.info(f"✅ Page content loaded successfully: {len(html_content)} characters")
            print(f"✅ HTML content fetched: {len(html_content)} chars")
            
        except Exception as request_error:
            logger.error(f"requests error: {request_error}")
            print(f"❌ Failed to fetch URL: {request_error}")
            return fallback
        
        # BeautifulSoupでHTMLを解析
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # テキストコンテンツを抽出（スクリプトやスタイルを除外）
        for script in soup(["script", "style"]):
            script.decompose()
        
        text_content = soup.get_text()
        # 空白を整理
        lines = (line.strip() for line in text_content.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)
        
        # テキストが長すぎる場合は切り詰め（OpenAI APIのトークン制限のため）
        if len(text) > 3000:
            text = text[:3000]
        
        logger.info(f"📝 Extracted text content: {len(text)} characters")
        
        # OpenAI APIで要約
        openai_api_key = os.environ.get('OPENAI_API_KEY', '')
        
        if not openai_api_key:
            logger.warning("OPENAI_API_KEY not set in environment variables")
            return "Website analysis unavailable (API key not configured)" if language == 'en' else "ウェブサイト分析が利用できません（APIキーが未設定）"
        
        prompt = f"""以下のウェブサイトの内容を分析し、以下の情報を含む簡潔な要約（150-200文字）を作成してください：
- 販売している商品・サービスの種類と特徴
- ビジネスのポジショニングや独自性
- 主な特徴や強み

ウェブサイトの内容:
{text}

要約は{'日本語' if language == 'ja' else '英語'}で作成してください。"""

        summary_response = requests.post(
            'https://api.openai.com/v1/chat/completions',
            headers={
                'Authorization': f'Bearer {openai_api_key}',
                'Content-Type': 'application/json'
            },
            json={
                'model': 'gpt-4o-mini',
                'messages': [
                    {'role': 'user', 'content': prompt}
                ],
                'max_tokens': 300,
                'temperature': 0.7
            },
            timeout=30
        )
        
        if summary_response.status_code == 200:
            result = summary_response.json()
            summary = result['choices'][0]['message']['content']
            logger.info(f"✅ Website summary generated: {summary[:100]}...")
            return summary
        else:
            logger.error(f"OpenAI API error: {summary_response.status_code}")
            logger.error(f"Response body: {summary_response.text[:500]}")
            return fallback
            
    except Exception as e:
        logger.error(f"Website crawl and analysis error: {e}")
        logger.error(traceback.format_exc())
        fallback = '手動でサイト概要を入力してください' if language == 'ja' else 'Please manually enter website description here'
        return fallback

@app.route('/api/create-pptx', methods=['POST'])
def create_pptx():
    """PPTXスライドを生成"""
    # 強制的にログ出力（デバッグ用）
    print("=" * 80)
    print("🚀 CREATE_PPTX FUNCTION CALLED - STARTING EXECUTION")
    print("=" * 80)
    try:
        data = request.json
        print(f"📦 Received data: {data}")
        channel_name = data.get('channel_name', '')
        industry = data.get('industry', '')
        country = data.get('country', '')
        url = data.get('url', '')
        language = data.get('language', 'ja')
        print(f"📝 Parsed: channel={channel_name}, industry={industry}, country={country}, url={url}, lang={language}")
        
        # 新しい指標を取得（CTA Click、50% Completion、100% Completion）
        cta_click_uu_rate_raw = data.get('cta_click_uu_rate', 'N/A')
        completion_50_uu_rate_raw = data.get('completion_50_uu_rate', 'N/A')
        completion_100_uu_rate_raw = data.get('completion_100_uu_rate', 'N/A')
        
        # パーセント表示用に100倍する（0.5 → 50%）
        def format_as_percentage(value):
            """数値を100倍してパーセント表示にする"""
            if value == 'N/A' or value is None or value == '':
                return 'N/A'
            try:
                numeric_value = float(value)
                percentage = numeric_value * 100
                # 小数点以下2桁まで表示
                return f"{percentage:.2f}%"
            except (ValueError, TypeError):
                return 'N/A'
        
        cta_click_uu_rate = format_as_percentage(cta_click_uu_rate_raw)
        completion_50_uu_rate = format_as_percentage(completion_50_uu_rate_raw)
        completion_100_uu_rate = format_as_percentage(completion_100_uu_rate_raw)
        
        print(f"✅ Percentage formatting complete: cta_click={cta_click_uu_rate}, completion_50={completion_50_uu_rate}, completion_100={completion_100_uu_rate}")
        
        logger.info(f"PPTX生成開始: Channel={channel_name}, 言語: {language}")
        logger.info(f"受信データ: channel_name={channel_name}, industry={industry}, country={country}, url={url}, format={data.get('format', 'NOT_PROVIDED')}")
        logger.info(f"指標データ(raw): CTA Click UU Rate={cta_click_uu_rate_raw}, 50% Completion UU Rate={completion_50_uu_rate_raw}, 100% Completion UU Rate={completion_100_uu_rate_raw}")
        logger.info(f"指標データ(formatted): CTA Click UU Rate={cta_click_uu_rate}, 50% Completion UU Rate={completion_50_uu_rate}, 100% Completion UU Rate={completion_100_uu_rate}")
        
        # ウェブサイト情報を抽出
        website_info = extract_website_info(url)
        
        # 言語が英語の場合、日本語テキストを翻訳
        if language == 'en':
            channel_name = translate_text(channel_name, 'en')
            company_details = translate_text(website_info['description'], 'en')
        else:
            company_details = website_info['description']
        
        # テンプレートを読み込む
        template_path = os.path.join(os.path.dirname(__file__), 'Template.pptx')
        prs = Presentation(template_path)
        
        # 言語に応じてスライドを選択（0: 日本語, 1: 英語）
        slide_index = 0 if language == 'ja' else 1
        slide = prs.slides[slide_index]
        
        # フォールバックメッセージ
        fallback_screenshot = '手動で画面キャプチャを貼り付けてください' if language == 'ja' else 'Please manually paste a screenshot here'
        fallback_logo = '手動でロゴを貼り付けてください' if language == 'ja' else 'Please manually paste the logo here'
        fallback_website = '手動でサイト概要を入力してください' if language == 'ja' else 'Please manually enter website description here'
        
        # WebクローラーとOpenAI APIでWebsite descriptionを取得
        website_description_enhanced = crawl_and_analyze_website(url, language) if url else fallback_website
        
        # 検出されたフォーマットを取得（dataから渡される場合、またはURLから検出）
        detected_format = data.get('format', 'Unknown')
        
        # fwタグが設置されているページのHTMLを取得して、Why firework?を生成
        fallback_why_firework = '目的: 動画で商品の魅力や使い方を分かりやすく説明 / 主要KPI: 視聴完了率' if language == 'ja' else 'Objective: Explain product features and usage through video / Key KPI: Video completion rate'
        why_firework_text = fallback_why_firework
        
        if url:
            try:
                logger.info(f"🔍 Starting Why firework generation for URL: {url}")
                # URLからHTMLコンテンツとフォーマットを取得
                has_fw, html_content, format_from_url = check_fw_tag_in_url(url)
                logger.info(f"📄 check_fw_tag_in_url result: has_fw={has_fw}, format_from_url={format_from_url}, html_length={len(html_content) if html_content else 0}")
                
                # フォーマットが不明な場合はURLから検出したものを使用
                if detected_format == 'Unknown' and format_from_url != 'Unknown':
                    detected_format = format_from_url
                    logger.info(f"✅ Format detected from URL: {detected_format}")
                
                if html_content:
                    logger.info(f"🚀 Calling generate_why_firework with format={detected_format}")
                    # website_description、Fireworkフォーマット情報も渡して、より正確な分析を行う
                    why_firework_text = generate_why_firework(url, html_content, website_description_enhanced, language, firework_format=detected_format)
                    logger.info(f"✅ Why firework text generated: {why_firework_text}")
                    
                    # 「目的: 」プレフィックスがない場合は追加（フォールバック）
                    if language == 'ja':
                        if not why_firework_text.startswith('目的:') and not why_firework_text.startswith('目的：'):
                            why_firework_text = '目的: ' + why_firework_text
                            logger.info(f"Added '目的: ' prefix: {why_firework_text}")
                    else:
                        if not why_firework_text.startswith('Objective:'):
                            why_firework_text = 'Objective: ' + why_firework_text
                            logger.info(f"Added 'Objective: ' prefix: {why_firework_text}")
                else:
                    logger.warning("HTML content not available for Why firework generation")
            except Exception as e:
                logger.error(f"Error generating Why firework: {e}")
                logger.error(traceback.format_exc())
        
        # プレースホルダーのテキストを置換（Business NameとCompany detailsは削除）
        replacements = {
            '{Business Country}': country,
            '{Account: Industry}': industry,
            '{Channel Name}': channel_name,
            '{URL}': url,
            '{Website description}': website_description_enhanced,
            '{Why firework?}': why_firework_text,
            '{Format}': detected_format,  # フォーマットを追加
            '{CTA Click UU Rate}': cta_click_uu_rate,  # パーセント表示済み
            '{50% Completion UU Rate}': completion_50_uu_rate,  # パーセント表示済み
            '{100% Completion UU Rate}': completion_100_uu_rate  # パーセント表示済み
        }
        
        # デバッグ: replacements辞書の内容をログ出力
        logger.info("========== Replacements Dictionary ==========")
        for key, value in replacements.items():
            value_preview = str(value)[:100] if value else 'None'
            logger.info(f"  {key}: {value_preview}")
        logger.info("==========================================")
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                original_text = shape.text
                new_text = original_text
                
                # すべてのプレースホルダーを置換
                for placeholder, value in replacements.items():
                    if placeholder in new_text:
                        logger.info(f"🔄 Found placeholder '{placeholder}' in shape. Replacing with: {str(value)[:50]}...")
                        new_text = new_text.replace(placeholder, value)
                
                # テキストが変更された場合のみ更新
                if new_text != original_text:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.text = new_text
                        
                        # {Website description}の場合、フォントサイズを10.5ptに設定
                        if '{Website description}' in original_text:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10.5)
                        
                        # {Why firework?}の場合、フォントサイズを12ptに設定し、disclaimerを赤字で追加
                        if '{Why firework?}' in original_text:
                            # disclaimerを追加
                            disclaimer_text = ' （想定内容につき、要加工）' if language == 'ja' else ' (Estimated content, editing required)'
                            
                            # 「」（かっこ）を削除
                            clean_text = why_firework_text.strip()
                            if clean_text.startswith('「') and clean_text.endswith('」'):
                                clean_text = clean_text[1:-1]
                            if clean_text.startswith('"') and clean_text.endswith('"'):
                                clean_text = clean_text[1:-1]
                            
                            # 「 / 主要KPI: 」を改行に変更（日本語）
                            if ' / 主要KPI: ' in clean_text:
                                clean_text = clean_text.replace(' / 主要KPI: ', '\n主要KPI: ')
                            # 「 / Key KPI: 」を改行に変更（英語）
                            if ' / Key KPI: ' in clean_text:
                                clean_text = clean_text.replace(' / Key KPI: ', '\nKey KPI: ')
                            
                            # テキストフレームをクリアして再構築
                            shape.text_frame.clear()
                            
                            # 本文を追加（黒字、12pt）
                            p = shape.text_frame.paragraphs[0]
                            run_main = p.add_run()
                            run_main.text = clean_text
                            run_main.font.size = Pt(12)
                            run_main.font.color.rgb = RGBColor(0, 0, 0)  # 黒字
                            
                            # disclaimerを追加（赤字、12pt）
                            run_disclaimer = p.add_run()
                            run_disclaimer.text = disclaimer_text
                            run_disclaimer.font.size = Pt(12)
                            run_disclaimer.font.color.rgb = RGBColor(255, 0, 0)  # 赤字
                    else:
                        shape.text = new_text
        
        # Playwrightを使用してスクリーンショットを取得して挿入
        screenshot_inserted = False
        
        # URLから検出されたフォーマット情報を取得（既にdataから取得済みの場合はそれを使用）
        if detected_format == 'Unknown' and url:
            try:
                has_fw, html_content, format_temp = check_fw_tag_in_url(url)
                detected_format = format_temp
                logger.info(f"Detected format for screenshot: {detected_format}")
            except Exception as e:
                logger.warning(f"Could not detect format: {e}")
        
        if url:
            try:
                logger.info(f"Generating screenshot for URL: {url}")
                print("📸📸📸 Starting screenshot capture for {Insert Screenshot here}...")
                
                # {Insert Screenshot here}: ページ上部のスクリーンショット（外部API使用）
                logger.info("Capturing top portion of page with external API...")
                print("🌐 Capturing top portion of page (external API)...")
                # 上部800pxのスクリーンショットを取得
                img_data = capture_screenshot_with_api(url, width=1200, height=800)
                
                if img_data:
                    # 画像サイズをチェック（白い画像を検出）
                    img_data.seek(0)
                    img_size = len(img_data.getvalue())
                    logger.info(f"Screenshot image size: {img_size} bytes")
                    print(f"📊 Screenshot size: {img_size} bytes")
                    
                    # 10KB未満の画像は白い画像とみなす
                    if img_size < 10000:
                        logger.warning(f"⚠️ Screenshot is too small ({img_size} bytes), likely a blank image. Will use HTML file instead.")
                        print(f"⚠️ Image too small: {img_size} bytes (threshold: 10000 bytes)")
                        img_data = None  # HTMLファイル保存に進む
                    else:
                        print(f"✅ Image size OK: {img_size} bytes, proceeding with insertion...")
                        # 画像を開いて確認
                        img_data.seek(0)
                        img = Image.open(img_data)
                        logger.info(f"Image dimensions: {img.width}x{img.height}")
                        print(f"📐 Image dimensions: {img.width}x{img.height}")
                        
                        # 画像を挿入する位置を探す
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and '{Insert Screenshot here}' in shape.text:
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height
                                
                                # プレースホルダーを削除
                                sp = shape.element
                                sp.getparent().remove(sp)
                                
                                # BytesIOを再度シーク（重要！）
                                img_data.seek(0)
                                
                                # 画像を枠のサイズに合わせて挿入
                                slide.shapes.add_picture(img_data, left, top, width=width, height=height)
                                screenshot_inserted = True
                                logger.info(f"✅ Screenshot inserted successfully ({img_size} bytes) at ({left}, {top}) with frame size {width}x{height}")
                                print(f"✅✅✅ Screenshot inserted into {{Insert Screenshot here}} successfully! Fitted to frame: {width}x{height}")
                                break
                else:
                    logger.warning(f"External API screenshot failed - no image data returned")
                    print(f"❌ External API returned no image data for {{Insert Screenshot here}}")
            except Exception as e:
                logger.warning(f"スクリーンショット取得失敗: {e}")
                logger.warning(traceback.format_exc())
        
        # スクリーンショットが挿入できなかった場合、HTMLファイルを保存してリンクを追加
        if not screenshot_inserted and url:
            # HTMLファイルを保存
            html_filename = f"{channel_name.replace(' ', '_')}_page.html"
            html_path = os.path.join(os.path.dirname(__file__), html_filename)
            
            if save_complete_html_page(url, html_path):
                logger.info(f"HTML page saved: {html_path}")
                
                # テキストボックスを探してHTMLリンクを追加
                for shape in slide.shapes:
                    if hasattr(shape, "text") and '{Insert Screenshot here}' in shape.text:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.clear()
                            p = shape.text_frame.paragraphs[0]
                            
                            # メインテキスト
                            run1 = p.add_run()
                            run1.text = "スクリーンショット取得失敗\n\n" if language == 'ja' else "Screenshot capture failed\n\n"
                            run1.font.size = Pt(12)
                            run1.font.color.rgb = RGBColor(255, 0, 0)
                            
                            # HTMLファイル説明
                            run2 = p.add_run()
                            run2.text = f"📄 代わりにHTMLファイルを保存しました:\n{html_filename}\n\n" if language == 'ja' else f"📄 HTML file saved instead:\n{html_filename}\n\n"
                            run2.font.size = Pt(10)
                            run2.font.color.rgb = RGBColor(0, 0, 0)
                            
                            # 使用方法
                            run3 = p.add_run()
                            run3.text = "使用方法: PPTXファイルと同じフォルダにHTMLファイルがあります。\nブラウザで開いてスクロールしてFireworkフォーマットを確認できます。" if language == 'ja' else "Usage: HTML file is in the same folder as PPTX.\nOpen in browser and scroll to view Firework format."
                            run3.font.size = Pt(9)
                            run3.font.color.rgb = RGBColor(100, 100, 100)
                        else:
                            shape.text = f"Screenshot failed. HTML saved: {html_filename}"
                        break
            else:
                # HTMLファイル保存も失敗した場合
                for shape in slide.shapes:
                    if hasattr(shape, "text") and '{Insert Screenshot here}' in shape.text:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.text = fallback_screenshot
                        else:
                            shape.text = fallback_screenshot
                        break
        
        # ロゴを検索して挿入（Template 3では3つ、改善された検索クエリを使用）
        logo_count = 3  # Template.pptxは3つのロゴプレースホルダー
        
        # Logo 1 & 2: 通常の検索クエリ
        logo_urls_12 = search_logo_images(channel_name, country=country, industry=industry, count=2)
        logger.info(f"Found {len(logo_urls_12)} logo URLs (1&2) for {channel_name} (country={country}, industry={industry})")
        
        # Logo 3: 異なるデザインを取得するため別のクエリで検索
        # "{Channel Name}+{Industry}+icon" または "alternative logo"で検索
        logo_urls_3 = search_logo_images(channel_name, country='', industry=industry + ' icon', count=3)
        logger.info(f"Found {len(logo_urls_3)} alternative logo URLs (3) for {channel_name}")
        
        # すべてのロゴURLを結合（Logo 3は異なる検索結果から）
        logo_urls = logo_urls_12[:2]  # Logo 1 & 2
        if logo_urls_3:
            # Logo 3: logo_urls_12と重複しないデザインを選択
            for url in logo_urls_3:
                if url not in logo_urls:
                    logo_urls.append(url)
                    break
            # まだ3つ目がない場合、logo_urls_3の最後を追加
            if len(logo_urls) < 3 and len(logo_urls_3) > 0:
                logo_urls.append(logo_urls_3[-1])
        
        # 各ロゴプレースホルダーを探して挿入
        for logo_index in range(1, logo_count + 1):  # 1, 2, 3
            placeholder = f'{{Channel logo {logo_index}}}'
            logo_inserted = False
            
            # 対応するロゴURLがある場合
            if logo_index <= len(logo_urls):
                logo_url = logo_urls[logo_index - 1]
                try:
                    logo_response = requests.get(logo_url, timeout=10)
                    if logo_response.status_code == 200:
                        logo_data = io.BytesIO(logo_response.content)
                        logo_img = Image.open(logo_data)
                        
                        # ロゴを挿入する位置を探す
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and placeholder in shape.text:
                                left = shape.left
                                top = shape.top
                                max_width = shape.width
                                max_height = shape.height
                                
                                # プレースホルダーを削除
                                sp = shape.element
                                sp.getparent().remove(sp)
                                
                                # アスペクト比を保持してリサイズ
                                img_width, img_height = logo_img.size
                                aspect = img_width / img_height
                                
                                if max_width / max_height > aspect:
                                    new_height = max_height
                                    new_width = int(max_height * aspect)
                                else:
                                    new_width = max_width
                                    new_height = int(max_width / aspect)
                                
                                # 画像を挿入
                                slide.shapes.add_picture(logo_data, left, top, width=new_width, height=new_height)
                                logo_inserted = True
                                logger.info(f"Logo {logo_index} inserted successfully")
                                break
                except Exception as e:
                    logger.warning(f"ロゴ{logo_index}取得失敗: {e}")
            
            # ロゴが挿入できなかった場合、フォールバックテキストを表示
            if not logo_inserted:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and placeholder in shape.text:
                        if hasattr(shape, "text_frame"):
                            shape.text_frame.text = fallback_logo
                        else:
                            shape.text = fallback_logo
                        logger.info(f"Logo {logo_index} fallback text inserted")
                        break
        
        # 選択したスライド以外を削除
        slides_to_delete = []
        for i, s in enumerate(prs.slides):
            if i != slide_index:
                slides_to_delete.append(i)
        
        # 逆順で削除（インデックスの変更を避けるため）
        for idx in reversed(slides_to_delete):
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
        
        # メモリ上にPPTXを保存
        pptx_io = io.BytesIO()
        try:
            prs.save(pptx_io)
            pptx_io.seek(0)
            
            # PPTXファイルサイズを検証
            pptx_size = len(pptx_io.getvalue())
            logger.info(f"PPTX生成完了: {channel_name}, サイズ: {pptx_size} bytes")
            
            # 最小サイズチェック（10KB未満は異常）
            if pptx_size < 10000:
                logger.error(f"❌ PPTX file too small ({pptx_size} bytes), likely corrupted!")
                raise Exception(f"Generated PPTX file is too small: {pptx_size} bytes")
            
            # 先頭に戻す
            pptx_io.seek(0)
            
        except Exception as pptx_save_error:
            logger.error(f"PPTX保存エラー: {pptx_save_error}")
            logger.error(traceback.format_exc())
            raise
        
        # HTMLファイルが存在する場合はZIPファイルとして返す
        html_filename = f"{channel_name.replace(' ', '_')}_page.html"
        html_path = os.path.join(os.path.dirname(__file__), html_filename)
        
        if os.path.exists(html_path):
            import zipfile
            
            # ZIPファイルを作成
            zip_io = io.BytesIO()
            # ファイル名をサニタイズ（スペースとカンマを削除）
            safe_filename = channel_name.replace(' ', '_').replace(',', '').replace('.', '')
            safe_html_filename = html_filename.replace(' ', '_').replace(',', '').replace('.', '')
            
            with zipfile.ZipFile(zip_io, 'w', zipfile.ZIP_DEFLATED) as zipf:
                # PPTXファイルを追加（バイナリデータをそのまま書き込む）
                pptx_data = pptx_io.getvalue()
                logger.info(f"Adding PPTX to ZIP: {len(pptx_data)} bytes")
                zipf.writestr(f'{safe_filename}_Casestudy.pptx', pptx_data)
                
                # HTMLファイルを追加
                with open(html_path, 'r', encoding='utf-8') as f:
                    zipf.writestr(safe_html_filename, f.read())
                
                # 説明ファイルを追加
                readme_text = """HTMLファイルの使用方法 / How to use HTML file
==============================================

1. このZIPファイルを解凍してください / Extract this ZIP file
2. PPTXファイルとHTMLファイルが含まれています / Contains PPTX and HTML files
3. HTMLファイルをブラウザで開いてください / Open HTML file in browser
4. スクロールしてFireworkフォーマットを確認できます / Scroll to view Firework format

注意: スクリーンショット自動取得に失敗したため、HTMLファイルを提供しています。
Note: HTML file provided because automatic screenshot capture failed.
"""
                zipf.writestr('README.txt', readme_text)
            
            # HTMLファイルを削除（クリーンアップ）
            try:
                os.remove(html_path)
            except:
                pass
            
            zip_io.seek(0)
            logger.info(f"ZIP file created with HTML: {channel_name}")
            
            # ファイル名をサニタイズ（スペースとカンマを削除）
            safe_filename = channel_name.replace(' ', '_').replace(',', '').replace('.', '')
            
            return send_file(
                zip_io,
                mimetype='application/zip',
                as_attachment=True,
                download_name=f'{safe_filename}_Casestudy.zip'
            )
        else:
            # HTMLファイルがない場合はPPTXのみ返す
            # ファイル名をサニタイズ（スペースとカンマを削除）
            safe_filename = channel_name.replace(' ', '_').replace(',', '').replace('.', '')
            
            return send_file(
                pptx_io,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f'{safe_filename}_Casestudy.pptx'
            )
    
    except Exception as e:
        logger.error(f"PPTX生成エラー: {e}")
        logger.error(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

@app.route('/api/export-excel', methods=['POST'])
def export_excel():
    """Excelファイルをエクスポート"""
    try:
        data = request.json
        rows = data.get('data', [])
        columns = data.get('columns', [])
        language = data.get('language', 'ja')
        
        logger.info(f"Excelエクスポート開始: {len(rows)}行, 言語: {language}")
        
        # 列名を言語に応じて変換（会社名とビジネス名は削除済み）
        if language == 'en':
            column_mapping = {
                'チャンネル名': 'Channel Name',
                '業種': 'Industry',
                '国': 'Country',
                'URL': 'URL'
            }
            translated_columns = [column_mapping.get(col, col) for col in columns]
        else:
            translated_columns = columns
        
        # DataFrameを作成
        df = pd.DataFrame(rows, columns=columns)
        df.columns = translated_columns
        
        # Excelファイルを作成
        excel_io = io.BytesIO()
        with pd.ExcelWriter(excel_io, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Results', index=False)
            
            # スタイルを適用
            workbook = writer.book
            worksheet = writer.sheets['Results']
            
            # ヘッダー行のスタイル
            header_fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
            header_font = Font(bold=True, color='FFFFFF')
            
            for cell in worksheet[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center', vertical='center')
            
            # 列幅を自動調整
            for column in worksheet.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                worksheet.column_dimensions[column_letter].width = adjusted_width
        
        excel_io.seek(0)
        
        logger.info(f"Excelエクスポート完了: {len(rows)}行")
        
        return send_file(
            excel_io,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name='analysis_result.xlsx'
        )
    
    except Exception as e:
        logger.error(f"Excelエクスポートエラー: {e}")
        logger.error(traceback.format_exc())
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    port = int(os.environ.get('PORT', 5000))
    debug_mode = os.environ.get('FLASK_ENV', 'production') != 'production'
    app.run(host='0.0.0.0', port=port, debug=debug_mode)
